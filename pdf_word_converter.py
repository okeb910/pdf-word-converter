"""
PDF ↔ Word/PowerPoint 互相转换工具

PDF → Word:
  - Windows Microsoft Word 原生转换（推荐，质量通常较好）
  - 整页渲染为图片嵌入 DOCX（视觉接近，文字一般不可编辑）
  - LibreOffice 无头转换（免费备选）

Word → PDF:
  - Windows 使用 Word COM，macOS 使用 Word AppleScript，均可回退 LibreOffice

PDF → PowerPoint:
  - 每页自适应高清渲染为整页图片，外观稳定但元素不可编辑

PowerPoint → PDF:
  - Windows 使用 PowerPoint COM，macOS 使用 PowerPoint AppleScript，均可回退 LibreOffice
"""
import hashlib
import io
import logging
import os
import subprocess
import sys
import tempfile
import threading
import time
import tkinter as tk
import unicodedata
from collections import Counter
from dataclasses import dataclass, field, replace
from difflib import SequenceMatcher
from pathlib import Path
from tkinter import filedialog, font as tkfont, messagebox, ttk

try:
    from tkinterdnd2 import COPY, DND_FILES, REFUSE_DROP, TkinterDnD
except Exception:  # Drag and drop is optional; standard Tk remains usable.
    COPY = "copy"
    DND_FILES = None
    REFUSE_DROP = "refuse_drop"
    TkinterDnD = None

from app_environment import (
    APP_VERSION,
    build_install_command,
    find_winget,
    open_official_download,
    run_install_command,
)
from engine_models import EngineState, EngineStatus
from conversion_specs import (
    PDF_TARGET_POWERPOINT,
    PDF_TARGET_WORD,
    resolve_conversion_spec,
)
from drop_logic import MixedSourceKindsError, classify_dropped_paths
from macos_office import MacOSOfficeError, MacOSPowerPointBackend, MacOSWordBackend
from pdf_fidelity import (
    PdfFidelityAnalysisCancelled,
    PdfFidelityRisk,
    analyze_pdf_fidelity_risk,
    inspect_editable_docx_tables,
)
from docx_table_repair import DocxTableRepairError, repair_docx_table_topology
from platform_services import InstallerActionKind, create_platform_services


PLATFORM_SERVICES = create_platform_services()
IS_MACOS = PLATFORM_SERVICES.platform == "darwin"
WORD_NATIVE = "word_native"
POWERPOINT_NATIVE = "powerpoint_native"
LEGACY_ENGINE_ALIASES = {
    "word_com": WORD_NATIVE,
    "powerpoint_com": POWERPOINT_NATIVE,
}

@dataclass(slots=True)
class PdfWordBatchPolicy:
    """Per-file PDF-to-Word fidelity choices for one serial batch."""

    default_method: str
    method_overrides: dict[Path, str] = field(default_factory=dict)
    editable_table_reports: dict[Path, PdfFidelityRisk] = field(default_factory=dict)
    selectable_text_reports: dict[Path, PdfFidelityRisk] = field(default_factory=dict)
    completion_warnings: dict[Path, str] = field(default_factory=dict)
    advisory_validation_paths: set[Path] = field(default_factory=set)
    blocked_paths: dict[Path, str] = field(default_factory=dict)
    unverified_paths: set[Path] = field(default_factory=set)
    image_protected_paths: set[Path] = field(default_factory=set)
    source_identities: dict[Path, tuple[int, int, int, int, str]] = field(
        default_factory=dict
    )
    source_snapshots: dict[Path, Path] = field(default_factory=dict)
    snapshot_directories: list[object] = field(default_factory=list, repr=False)

    def cleanup_snapshots(self) -> None:
        while self.snapshot_directories:
            temporary_directory = self.snapshot_directories.pop()
            try:
                temporary_directory.cleanup()
            except OSError:
                pass
        self.source_snapshots.clear()


def _source_file_identity(path) -> tuple[int, int, int, int, str]:
    source_path = Path(path)
    stat_before = source_path.stat()
    digest = hashlib.sha256()
    with source_path.open("rb") as source_stream:
        while chunk := source_stream.read(1024 * 1024):
            digest.update(chunk)
    stat_after = source_path.stat()
    metadata_before = (
        int(getattr(stat_before, "st_dev", 0) or 0),
        int(getattr(stat_before, "st_ino", 0) or 0),
        int(stat_before.st_size),
        int(stat_before.st_mtime_ns),
    )
    metadata_after = (
        int(getattr(stat_after, "st_dev", 0) or 0),
        int(getattr(stat_after, "st_ino", 0) or 0),
        int(stat_after.st_size),
        int(stat_after.st_mtime_ns),
    )
    if metadata_before != metadata_after:
        raise RuntimeError("读取源 PDF 期间文件发生变化")
    return (*metadata_after, digest.hexdigest())


def _copy_verified_source_snapshot(source, destination, expected_identity):
    """Copy the exact preflight bytes to a private conversion snapshot."""

    source_path = Path(source)
    destination_path = Path(destination)
    digest = hashlib.sha256()
    destination_path.parent.mkdir(parents=True, exist_ok=True)
    with source_path.open("rb") as source_stream, destination_path.open("xb") as output:
        stat_before = os.fstat(source_stream.fileno())
        while chunk := source_stream.read(1024 * 1024):
            digest.update(chunk)
            output.write(chunk)
        output.flush()
        os.fsync(output.fileno())
        stat_after = os.fstat(source_stream.fileno())
    metadata_before = (
        int(getattr(stat_before, "st_dev", 0) or 0),
        int(getattr(stat_before, "st_ino", 0) or 0),
        int(stat_before.st_size),
        int(stat_before.st_mtime_ns),
    )
    metadata_after = (
        int(getattr(stat_after, "st_dev", 0) or 0),
        int(getattr(stat_after, "st_ino", 0) or 0),
        int(stat_after.st_size),
        int(stat_after.st_mtime_ns),
    )
    copied_identity = (*metadata_after, digest.hexdigest())
    if metadata_before != metadata_after or copied_identity != expected_identity:
        destination_path.unlink(missing_ok=True)
        raise RuntimeError("源 PDF 在创建转换快照期间发生变化，请重新加入后再转换")
    return destination_path


def _normalize_table_text(value) -> str:
    characters = []
    for character in unicodedata.normalize("NFC", str(value or "")):
        if unicodedata.category(character).startswith("P"):
            compatible = unicodedata.normalize("NFKC", character)
            if (
                len(compatible) == 1
                and unicodedata.category(compatible).startswith("P")
            ):
                character = compatible
        characters.append(character)

    normalized = []
    for index, character in enumerate(characters):
        if character in "\r\n\v\f":
            continue
        if not character.isspace():
            normalized.append(character)
            continue

        previous = next(
            (
                value
                for value in reversed(normalized)
                if not value.isspace()
            ),
            "",
        )
        following = next(
            (
                value
                for value in characters[index + 1 :]
                if not value.isspace()
            ),
            "",
        )
        preserve_space = bool(
            previous
            and following
            and previous.isascii()
            and following.isascii()
            and previous.isalnum()
            and following.isalnum()
        )
        if preserve_space and (not normalized or normalized[-1] != " "):
            normalized.append(" ")
    return "".join(normalized).strip()


def _document_text_evidence_matches(
    source_text,
    output_text,
    *,
    source_character_count=None,
    output_character_count=None,
) -> bool:
    if source_text is None or output_text is None:
        return False
    source_text = str(source_text)
    output_text = str(output_text)
    measured_source_count = sum(
        not character.isspace() for character in source_text
    )
    measured_output_count = sum(
        not character.isspace() for character in output_text
    )
    if (
        source_character_count is not None
        and measured_source_count != int(source_character_count)
    ):
        return False
    if (
        output_character_count is not None
        and measured_output_count != int(output_character_count)
    ):
        return False
    normalized_source = _normalize_table_text(source_text)
    normalized_output = _normalize_table_text(output_text)
    return bool(
        len(normalized_source) == len(normalized_output)
        and Counter(normalized_source) == Counter(normalized_output)
    )


def _ordered_text_overlap_ratio(source_text, output_text) -> float:
    if not source_text:
        return 1.0
    if len(source_text) <= 20000 and len(output_text) <= 50000:
        matching_characters = sum(
            block.size
            for block in SequenceMatcher(
                None,
                source_text,
                output_text,
                autojunk=False,
            ).get_matching_blocks()
        )
        return matching_characters / len(source_text)

    gram_size = 3 if len(source_text) >= 3 else 1
    source_grams = Counter(
        source_text[index : index + gram_size]
        for index in range(len(source_text) - gram_size + 1)
    )
    output_grams = Counter(
        output_text[index : index + gram_size]
        for index in range(len(output_text) - gram_size + 1)
    )
    matching_grams = sum((source_grams & output_grams).values())
    return matching_grams / max(1, sum(source_grams.values()))


def _table_texts_preserved(source_texts, output_texts, minimum_ratio=0.8):
    sources = [
        text
        for text in (_normalize_table_text(value) for value in source_texts)
        if text
    ]
    outputs = {
        index: text
        for index, text in enumerate(
            _normalize_table_text(value) for value in output_texts
        )
        if text
    }
    matched = 0
    for source_text in sorted(sources, key=len, reverse=True):
        source_characters = Counter(source_text)
        best_index = None
        best_score = -1.0
        for index, output_text in outputs.items():
            character_overlap = sum(
                (source_characters & Counter(output_text)).values()
            ) / len(source_text)
            ordered_overlap = _ordered_text_overlap_ratio(
                source_text,
                output_text,
            )
            score = min(character_overlap, ordered_overlap)
            if score > best_score:
                best_index = index
                best_score = score
        if best_index is None or best_score < float(minimum_ratio):
            return False, matched, len(sources)
        outputs.pop(best_index)
        matched += 1
    return True, matched, len(sources)


def _coerce_table_shape(value):
    try:
        rows, columns = value
        rows = int(rows)
        columns = int(columns)
    except (TypeError, ValueError):
        return None
    if rows < 1 or columns < 1:
        return None
    return rows, columns


def _coerce_table_cell_count(value):
    try:
        count = int(value)
    except (TypeError, ValueError):
        return None
    return count if count > 0 else None


def _coerce_table_cell_spans(value, shape):
    shape = _coerce_table_shape(shape)
    if shape is None:
        return None
    try:
        spans = tuple(
            tuple(int(component) for component in span)
            for span in value
        )
    except (TypeError, ValueError):
        return None
    rows, columns = shape
    coverage = [[False] * columns for _row in range(rows)]
    for span in spans:
        if len(span) != 4:
            return None
        row, column, row_span, column_span = span
        if (
            row < 0
            or column < 0
            or row_span < 1
            or column_span < 1
            or row + row_span > rows
            or column + column_span > columns
        ):
            return None
        for covered_row in range(row, row + row_span):
            for covered_column in range(column, column + column_span):
                if coverage[covered_row][covered_column]:
                    return None
                coverage[covered_row][covered_column] = True
    if not spans or not all(all(row) for row in coverage):
        return None
    return tuple(sorted(spans))


def _coerce_table_cell_border_map(value, spans, shape):
    normalized_spans = _coerce_table_cell_spans(spans, shape)
    if normalized_spans is None:
        return None
    try:
        ordered_spans = tuple(
            tuple(int(component) for component in span)
            for span in spans
        )
        border_edges = tuple(tuple(edges) for edges in value)
    except (TypeError, ValueError):
        return None
    if len(ordered_spans) != len(normalized_spans) or len(border_edges) != len(
        ordered_spans
    ):
        return None
    if any(
        len(edges) != 4 or any(type(edge) is not bool for edge in edges)
        for edges in border_edges
    ):
        return None
    return tuple(
        sorted(
            zip(ordered_spans, border_edges),
            key=lambda item: item[0],
        )
    )


def _coerce_table_cell_matrix(value, shape):
    shape = _coerce_table_shape(shape)
    if shape is None:
        return None
    try:
        rows = tuple(tuple(row) for row in value)
    except TypeError:
        return None
    expected_rows, expected_columns = shape
    if (
        len(rows) != expected_rows
        or any(len(row) != expected_columns for row in rows)
    ):
        return None
    return tuple(
        tuple(None if cell is None else str(cell) for cell in row)
        for row in rows
    )


def _table_cell_matrix_text(matrix) -> str:
    return "".join(
        str(cell or "")
        for row in matrix
        for cell in row
        if cell is not None
    )


def _table_cell_matrix_matches_text(matrix, shape, table_text) -> bool:
    coerced = _coerce_table_cell_matrix(matrix, shape)
    return bool(
        coerced is not None
        and _normalize_table_text(_table_cell_matrix_text(coerced))
        == _normalize_table_text(table_text)
    )


def _table_cell_matrix_preserved(
    source_matrix,
    output_matrix,
    source_shape,
    output_shape,
    source_full_width_span_rows=None,
) -> bool:
    del source_full_width_span_rows
    source = _coerce_table_cell_matrix(source_matrix, source_shape)
    output = _coerce_table_cell_matrix(output_matrix, output_shape)
    if source is None or output is None or source_shape != output_shape:
        return False
    for source_row, output_row in zip(source, output):
        for source_value, output_value in zip(source_row, output_row):
            if (source_value is None) != (output_value is None):
                return False
            if source_value is not None and (
                _normalize_table_text(source_value)
                != _normalize_table_text(output_value)
            ):
                return False
    return True
def _table_pair_preserved(
    source_text,
    source_shape,
    source_cell_count,
    output_text,
    output_shape,
    output_cell_count,
    minimum_ratio=0.98,
    source_cell_matrix=None,
    output_cell_matrix=None,
    source_full_width_span_rows=None,
    source_cell_spans=None,
    output_cell_spans=None,
):
    source_text = _normalize_table_text(source_text)
    output_text = _normalize_table_text(output_text)
    source_shape = _coerce_table_shape(source_shape)
    output_shape = _coerce_table_shape(output_shape)
    source_cell_count = _coerce_table_cell_count(source_cell_count)
    output_cell_count = _coerce_table_cell_count(output_cell_count)
    if (
        source_shape is None
        or output_shape is None
        or source_cell_count is None
        or output_cell_count is None
        or source_shape != output_shape
        or output_cell_count != source_cell_count
    ):
        return False

    if (source_cell_spans is None) != (output_cell_spans is None):
        return False
    if source_cell_spans is not None:
        source_spans = _coerce_table_cell_spans(
            source_cell_spans,
            source_shape,
        )
        output_spans = _coerce_table_cell_spans(
            output_cell_spans,
            output_shape,
        )
        if source_spans is None or source_spans != output_spans:
            return False

    if (source_cell_matrix is None) != (output_cell_matrix is None):
        return False
    matrix_verified = False
    if source_cell_matrix is not None:
        if not _table_cell_matrix_preserved(
            source_cell_matrix,
            output_cell_matrix,
            source_shape,
            output_shape,
            source_full_width_span_rows=source_full_width_span_rows,
        ):
            return False
        matrix_verified = True

    # A confirmed blank grid has no text to compare, but its exact row, column,
    # cell and merge matrices still prove that it remains an editable table.
    if not source_text or not output_text:
        return bool(matrix_verified and not source_text and not output_text)

    source_characters = Counter(source_text)
    character_overlap = sum(
        (source_characters & Counter(output_text)).values()
    ) / len(source_text)
    ordered_overlap = _ordered_text_overlap_ratio(source_text, output_text)
    text_size_similarity = min(len(source_text), len(output_text)) / max(
        len(source_text), len(output_text)
    )
    return min(
        character_overlap,
        ordered_overlap,
        text_size_similarity,
    ) >= max(0.98, float(minimum_ratio))

def _ordered_source_residual(source_text, matched_output_text):
    """Return source character positions not covered by the matched main table."""

    source_text = _normalize_table_text(source_text)
    matched_output_text = _normalize_table_text(matched_output_text)
    if len(source_text) > 20000 or len(matched_output_text) > 50000:
        return None

    covered = bytearray(len(source_text))
    matcher = SequenceMatcher(
        None,
        source_text,
        matched_output_text,
        autojunk=False,
    )
    for block in matcher.get_matching_blocks():
        for index in range(block.a, block.a + block.size):
            covered[index] = 1
    return "".join(
        character
        for index, character in enumerate(source_text)
        if not covered[index]
    )


def _table_structure_gap(source_shape, source_cell_count, output_shape, output_cell_count):
    source_shape = _coerce_table_shape(source_shape)
    output_shape = _coerce_table_shape(output_shape)
    source_cell_count = _coerce_table_cell_count(source_cell_count)
    output_cell_count = _coerce_table_cell_count(output_cell_count)
    if (
        source_shape is None
        or output_shape is None
        or source_cell_count is None
        or output_cell_count is None
    ):
        return float("inf")

    source_rows, source_columns = source_shape
    output_rows, output_columns = output_shape
    source_area = source_rows * source_columns
    output_area = output_rows * output_columns
    return (
        abs(source_rows - output_rows) / source_rows
        + abs(source_columns - output_columns) / source_columns
        + abs(source_area - output_area) / source_area
        + abs(source_cell_count - output_cell_count) / source_cell_count
    )


def _fragment_reduces_structure_gap(
    source_shape,
    source_cell_count,
    matched_shape,
    matched_cell_count,
    fragment_shape,
    fragment_cell_count,
):
    matched_shape = _coerce_table_shape(matched_shape)
    fragment_shape = _coerce_table_shape(fragment_shape)
    matched_cell_count = _coerce_table_cell_count(matched_cell_count)
    fragment_cell_count = _coerce_table_cell_count(fragment_cell_count)
    if (
        matched_shape is None
        or fragment_shape is None
        or matched_cell_count is None
        or fragment_cell_count is None
    ):
        return False

    before = _table_structure_gap(
        source_shape,
        source_cell_count,
        matched_shape,
        matched_cell_count,
    )
    vertical_shape = (
        matched_shape[0] + fragment_shape[0],
        max(matched_shape[1], fragment_shape[1]),
    )
    horizontal_shape = (
        max(matched_shape[0], fragment_shape[0]),
        matched_shape[1] + fragment_shape[1],
    )
    combined_cells = matched_cell_count + fragment_cell_count
    after = min(
        _table_structure_gap(
            source_shape,
            source_cell_count,
            vertical_shape,
            combined_cells,
        ),
        _table_structure_gap(
            source_shape,
            source_cell_count,
            horizontal_shape,
            combined_cells,
        ),
    )
    return after < before


def _fragment_residual_positions(residual, candidate):
    residual = _normalize_table_text(residual)
    candidate = _normalize_table_text(candidate)
    if not residual or not candidate:
        return frozenset()
    if len(residual) > 20000 or len(candidate) > 50000:
        return frozenset(range(len(residual)))

    matcher = SequenceMatcher(None, residual, candidate, autojunk=False)
    positions = {
        index
        for block in matcher.get_matching_blocks()
        for index in range(block.a, block.a + block.size)
    }
    character_hits = sum(
        (Counter(residual) & Counter(candidate)).values()
    )
    membership = min(character_hits, len(positions)) / len(candidate)
    required_membership = 1.0 if len(candidate) <= 3 else 0.8
    return frozenset(positions) if membership >= required_membership else frozenset()

def _unmatched_table_looks_like_source_fragment(
    source_text,
    source_shape,
    source_cell_count,
    matched_output_text,
    matched_output_shape,
    matched_output_cell_count,
    unmatched_output_text,
    unmatched_output_shape,
    unmatched_output_cell_count,
):
    residual = _ordered_source_residual(source_text, matched_output_text)
    if residual is None:
        return bool(_normalize_table_text(unmatched_output_text))
    if not residual:
        return False

    candidate = _normalize_table_text(unmatched_output_text)
    if not candidate:
        return False
    residual_positions = _fragment_residual_positions(residual, candidate)
    if not residual_positions:
        return False
    if len(residual_positions) >= len(residual):
        return True
    if len(candidate) > 3:
        return True
    return _fragment_reduces_structure_gap(
        source_shape,
        source_cell_count,
        matched_output_shape,
        matched_output_cell_count,
        unmatched_output_shape,
        unmatched_output_cell_count,
    )

def _table_structures_preserved(
    source_texts,
    source_shapes,
    source_cell_counts,
    output_texts,
    output_shapes,
    output_cell_counts,
    source_cell_matrices=None,
    output_cell_matrices=None,
    source_full_width_span_rows=None,
    source_cell_spans=None,
    output_cell_spans=None,
    minimum_ratio=0.98,
):
    """Match source tables to complete editable Word tables in reading order."""

    source_texts = tuple(source_texts or ())
    source_shapes = tuple(source_shapes or ())
    source_cell_counts = tuple(source_cell_counts or ())
    output_texts = tuple(output_texts or ())
    output_shapes = tuple(output_shapes or ())
    output_cell_counts = tuple(output_cell_counts or ())
    matrices_requested = (
        source_cell_matrices is not None
        or output_cell_matrices is not None
    )
    if matrices_requested and (
        source_cell_matrices is None
        or output_cell_matrices is None
    ):
        return False, 0, len(source_texts)
    source_cell_matrices = (
        tuple(source_cell_matrices or ())
        if matrices_requested
        else ()
    )
    output_cell_matrices = (
        tuple(output_cell_matrices or ())
        if matrices_requested
        else ()
    )
    spans_requested = source_cell_spans is not None or output_cell_spans is not None
    if spans_requested and (
        source_cell_spans is None or output_cell_spans is None
    ):
        return False, 0, len(source_texts)
    source_cell_spans = (
        tuple(source_cell_spans or ()) if spans_requested else ()
    )
    output_cell_spans = (
        tuple(output_cell_spans or ()) if spans_requested else ()
    )
    source_count = len(source_texts)
    if source_full_width_span_rows is None:
        source_full_width_span_rows = tuple(
            (False,) * int(shape[0])
            for shape in source_shapes
            if _coerce_table_shape(shape) is not None
        )
    else:
        try:
            source_full_width_span_rows = tuple(
                tuple(bool(value) for value in table_rows)
                for table_rows in source_full_width_span_rows
            )
        except TypeError:
            return False, 0, source_count
    if (
        not source_count
        or len(source_shapes) != source_count
        or len(source_cell_counts) != source_count
        or len(output_texts) != len(output_shapes)
        or len(output_texts) != len(output_cell_counts)
        or len(source_full_width_span_rows) != source_count
        or any(
            len(table_rows) != int(shape[0])
            for table_rows, shape in zip(
                source_full_width_span_rows,
                source_shapes,
            )
            if _coerce_table_shape(shape) is not None
        )
        or (
            matrices_requested
            and (
                len(source_cell_matrices) != source_count
                or len(output_cell_matrices) != len(output_texts)
            )
        )
        or (
            spans_requested
            and (
                len(source_cell_spans) != source_count
                or len(output_cell_spans) != len(output_texts)
            )
        )
    ):
        return False, 0, source_count

    best_paths = [()] * (len(output_texts) + 1)
    for source_index, (source_text, source_shape, source_cell_count) in enumerate(
        zip(source_texts, source_shapes, source_cell_counts),
        start=1,
    ):
        next_paths = best_paths.copy()
        for output_index, (output_text, output_shape, output_cell_count) in enumerate(
            zip(output_texts, output_shapes, output_cell_counts),
            start=1,
        ):
            if len(next_paths[output_index - 1]) > len(next_paths[output_index]):
                next_paths[output_index] = next_paths[output_index - 1]
            if (
                len(best_paths[output_index - 1]) == source_index - 1
                and _table_pair_preserved(
                    source_text,
                    source_shape,
                    source_cell_count,
                    output_text,
                    output_shape,
                    output_cell_count,
                    minimum_ratio=minimum_ratio,
                    source_cell_matrix=(
                        source_cell_matrices[source_index - 1]
                        if matrices_requested
                        else None
                    ),
                    output_cell_matrix=(
                        output_cell_matrices[output_index - 1]
                        if matrices_requested
                        else None
                    ),
                    source_full_width_span_rows=(
                        source_full_width_span_rows[source_index - 1]
                    ),
                    source_cell_spans=(
                        source_cell_spans[source_index - 1]
                        if spans_requested
                        else None
                    ),
                    output_cell_spans=(
                        output_cell_spans[output_index - 1]
                        if spans_requested
                        else None
                    ),
                )
            ):
                candidate = best_paths[output_index - 1] + (output_index - 1,)
                if len(candidate) > len(next_paths[output_index]):
                    next_paths[output_index] = candidate
        best_paths = next_paths

    best_path = max(best_paths, key=len, default=())
    matched = len(best_path)
    if matched != source_count:
        return False, matched, source_count

    matched_outputs = set(best_path)
    normalized_outputs = tuple(_normalize_table_text(text) for text in output_texts)
    unmatched_output_indices = tuple(
        output_index
        for output_index, output_text in enumerate(normalized_outputs)
        if output_index not in matched_outputs and output_text
    )
    for source_index, matched_output_index in enumerate(best_path):
        residual = _ordered_source_residual(
            source_texts[source_index],
            output_texts[matched_output_index],
        )
        aggregate_residual_positions = set()
        for unmatched_output_index in unmatched_output_indices:
            if residual is not None:
                aggregate_residual_positions.update(
                    _fragment_residual_positions(
                        residual,
                        output_texts[unmatched_output_index],
                    )
                )
            if _unmatched_table_looks_like_source_fragment(
                source_texts[source_index],
                source_shapes[source_index],
                source_cell_counts[source_index],
                output_texts[matched_output_index],
                output_shapes[matched_output_index],
                output_cell_counts[matched_output_index],
                output_texts[unmatched_output_index],
                output_shapes[unmatched_output_index],
                output_cell_counts[unmatched_output_index],
            ):
                return False, max(0, matched - 1), source_count
        if residual and len(aggregate_residual_positions) >= len(residual):
            return False, max(0, matched - 1), source_count
    return True, matched, source_count

def _canonical_engine_key(key: str) -> str:
    return LEGACY_ENGINE_ALIASES.get(key, key)


def _coerce_engine_status(result, error=None) -> EngineStatus:
    if isinstance(result, EngineStatus):
        return result
    if error:
        return EngineStatus(EngineState.LAUNCH_FAILED, str(error))
    return EngineStatus(
        EngineState.AVAILABLE if bool(result) else EngineState.MISSING,
    )


def _status_selectable(status: EngineStatus) -> bool:
    return status.state in {EngineState.AVAILABLE, EngineState.UNVERIFIED}


def _availability_from_result(result, error=None):
    """Map a probe result to the tri-state value consumed by the UI."""
    status = _coerce_engine_status(result, error)
    if status.state is EngineState.CHECKING:
        return None
    return _status_selectable(status)


def _engine_result_selectable(result) -> bool:
    """Return whether a cached/probed engine result may be used for conversion."""
    return _status_selectable(_coerce_engine_status(result))


# ═══════════════════════════════════════════════════════════
#  引擎可用性检测
# ═══════════════════════════════════════════════════════════

def _check_word_com_available():
    """Detect Word registration without launching Office during startup."""
    if sys.platform != "win32":
        return EngineStatus(EngineState.UNSUPPORTED, "当前平台不支持 Microsoft Word COM")
    if not _com_application_registered("Word.Application"):
        return EngineStatus(EngineState.MISSING, "未检测到 Microsoft Word COM 注册")
    return EngineStatus(
        EngineState.UNVERIFIED,
        "已检测到 Microsoft Word，首次转换时验证 COM",
    )


def _check_powerpoint_com_available():
    """Detect PowerPoint registration without launching Office during startup."""
    if sys.platform != "win32":
        return EngineStatus(EngineState.UNSUPPORTED, "当前平台不支持 Microsoft PowerPoint COM")
    if not _com_application_registered("PowerPoint.Application"):
        return EngineStatus(EngineState.MISSING, "未检测到 Microsoft PowerPoint COM 注册")
    return EngineStatus(
        EngineState.UNVERIFIED,
        "已检测到 Microsoft PowerPoint，首次转换时验证 COM",
    )



def _com_application_registered(prog_id: str) -> bool:
    """Check COM registration without starting the Office application."""
    if sys.platform != "win32":
        return False
    try:
        import winreg

        with winreg.OpenKey(winreg.HKEY_CLASSES_ROOT, rf"{prog_id}\CLSID") as key:
            clsid, _value_type = winreg.QueryValueEx(key, None)
        return bool(clsid)
    except (ImportError, OSError):
        return False

def _find_libreoffice():
    """查找当前平台的 LibreOffice 可执行文件路径。"""
    path = PLATFORM_SERVICES.find_libreoffice()
    return str(path) if path else None

_LO_PATH = None  # 缓存找到的 LibreOffice 路径

def _get_lo_path():
    global _LO_PATH
    if _LO_PATH is None:
        _LO_PATH = _find_libreoffice()
    return _LO_PATH

def _check_libreoffice_available():
    """Detect LibreOffice without launching it during application startup."""
    lo = _get_lo_path()
    if not lo:
        return EngineStatus(EngineState.MISSING, "未检测到 LibreOffice 可执行文件")
    return EngineStatus(
        EngineState.UNVERIFIED,
        f"已检测到 LibreOffice：{lo}；首次转换时验证",
    )


def _load_pymupdf():
    """兼容 PyMuPDF 新旧模块名，并确认核心 API 存在。"""
    try:
        import pymupdf
        module = pymupdf
    except ImportError:
        import fitz
        module = fitz
    if not hasattr(module, "open"):
        raise ImportError("PyMuPDF 模块缺少 open()，请重新安装 PyMuPDF")
    return module


def _check_pymupdf_available():
    """检测 PyMuPDF 是否可用。"""
    try:
        _load_pymupdf()
        return True
    except ImportError:
        return False


def _check_pptx_component_available():
    """检测内置 PPTX 写入组件是否可用。"""
    try:
        from pptx import Presentation
        return callable(Presentation)
    except ImportError:
        return False


# 每个引擎独立缓存和锁，允许启动检测真正并行执行。
_WORD_LOCK = threading.Lock()
_POWERPOINT_LOCK = threading.Lock()
_LIBREOFFICE_LOCK = threading.Lock()
_PYMUPDF_LOCK = threading.Lock()
_PPTX_COMPONENT_LOCK = threading.Lock()
_ENGINE_CACHE_GENERATION = 0
_WORD_CACHE_GENERATION = -1
_POWERPOINT_CACHE_GENERATION = -1
_LIBREOFFICE_CACHE_GENERATION = -1
_PYMUPDF_CACHE_GENERATION = -1
_PPTX_COMPONENT_CACHE_GENERATION = -1
_WORD_AVAILABLE = None
_POWERPOINT_AVAILABLE = None
_LIBREOFFICE_AVAILABLE = None
_PYMUPDF_AVAILABLE = None
_PPTX_COMPONENT_AVAILABLE = None


def reset_engine_cache():
    """Invalidate cached checks without waiting for an older probe to finish."""
    global _ENGINE_CACHE_GENERATION, _LO_PATH
    global _WORD_AVAILABLE, _POWERPOINT_AVAILABLE, _LIBREOFFICE_AVAILABLE
    global _PYMUPDF_AVAILABLE, _PPTX_COMPONENT_AVAILABLE
    _ENGINE_CACHE_GENERATION += 1
    _LO_PATH = None
    _WORD_AVAILABLE = None
    _POWERPOINT_AVAILABLE = None
    _LIBREOFFICE_AVAILABLE = None
    _PYMUPDF_AVAILABLE = None
    _PPTX_COMPONENT_AVAILABLE = None


def word_com_available():
    global _WORD_AVAILABLE, _WORD_CACHE_GENERATION
    generation = _ENGINE_CACHE_GENERATION
    with _WORD_LOCK:
        if _WORD_CACHE_GENERATION != generation:
            result = _check_word_com_available()
            if generation == _ENGINE_CACHE_GENERATION:
                _WORD_AVAILABLE = result
                _WORD_CACHE_GENERATION = generation
            return result
        return _WORD_AVAILABLE


def powerpoint_com_available():
    global _POWERPOINT_AVAILABLE, _POWERPOINT_CACHE_GENERATION
    generation = _ENGINE_CACHE_GENERATION
    with _POWERPOINT_LOCK:
        if _POWERPOINT_CACHE_GENERATION != generation:
            result = _check_powerpoint_com_available()
            if generation == _ENGINE_CACHE_GENERATION:
                _POWERPOINT_AVAILABLE = result
                _POWERPOINT_CACHE_GENERATION = generation
            return result
        return _POWERPOINT_AVAILABLE


def libreoffice_available():
    global _LIBREOFFICE_AVAILABLE, _LIBREOFFICE_CACHE_GENERATION
    generation = _ENGINE_CACHE_GENERATION
    with _LIBREOFFICE_LOCK:
        if _LIBREOFFICE_CACHE_GENERATION != generation:
            result = _check_libreoffice_available()
            if generation == _ENGINE_CACHE_GENERATION:
                _LIBREOFFICE_AVAILABLE = result
                _LIBREOFFICE_CACHE_GENERATION = generation
            return result
        return _LIBREOFFICE_AVAILABLE


def pymupdf_available():
    global _PYMUPDF_AVAILABLE, _PYMUPDF_CACHE_GENERATION
    generation = _ENGINE_CACHE_GENERATION
    with _PYMUPDF_LOCK:
        if _PYMUPDF_CACHE_GENERATION != generation:
            result = _check_pymupdf_available()
            if generation == _ENGINE_CACHE_GENERATION:
                _PYMUPDF_AVAILABLE = result
                _PYMUPDF_CACHE_GENERATION = generation
            return result
        return _PYMUPDF_AVAILABLE


def pptx_component_available():
    global _PPTX_COMPONENT_AVAILABLE, _PPTX_COMPONENT_CACHE_GENERATION
    generation = _ENGINE_CACHE_GENERATION
    with _PPTX_COMPONENT_LOCK:
        if _PPTX_COMPONENT_CACHE_GENERATION != generation:
            result = _check_pptx_component_available()
            if generation == _ENGINE_CACHE_GENERATION:
                _PPTX_COMPONENT_AVAILABLE = result
                _PPTX_COMPONENT_CACHE_GENERATION = generation
            return result
        return _PPTX_COMPONENT_AVAILABLE


# ═══════════════════════════════════════════════════════════
#  PDF → Word 转换方法
# ═══════════════════════════════════════════════════════════

def pdf_to_word_via_word(pdf_path: str, docx_path: str, progress) -> None:
    """
    [推荐] Microsoft Word 原生转换
    Word 2013+ 可打开 PDF 并另存为可编辑 DOCX。
    质量通常优于纯第三方解析，但仍受 PDF 结构与本机 Word 版本影响。
    """
    import comtypes
    import comtypes.client
    import ctypes
    import time

    progress("正在启动 Microsoft Word...", 5)
    comtypes.CoInitialize()
    try:
        word = comtypes.client.CreateObject("Word.Application", dynamic=True)
    except Exception:
        comtypes.CoUninitialize()
        raise
    word.Visible = False
    word.DisplayAlerts = 0  # wdAlertsNone
    doc = None

    # 后台线程：自动关闭 Word 可能弹出的 PDF 转换进度对话框
    stop_dismisser = False

    def dismiss_word_dialogs():
        """每隔 0.5 秒查找并关闭 Word 的转换进度弹窗"""
        user32 = ctypes.windll.user32
        WNDENUMPROC = ctypes.WINFUNCTYPE(
            ctypes.c_bool, ctypes.c_void_p, ctypes.c_void_p,
        )

        triggers = [
            "正在转换", "Converting", "Opening PDF",
            "PDF 正在", "Word 正在",
        ]

        def enum_proc(hwnd, lParam):
            if not user32.IsWindowVisible(hwnd):
                return True
            length = user32.GetWindowTextLengthW(hwnd)
            if length == 0:
                return True
            buf = ctypes.create_unicode_buffer(length + 1)
            user32.GetWindowTextW(hwnd, buf, length + 1)
            title = buf.value
            if any(t in title for t in triggers):
                user32.PostMessageW(hwnd, 0x0100, 0x0D, 0)
                user32.PostMessageW(hwnd, 0x0101, 0x0D, 0)
            return True

        callback = WNDENUMPROC(enum_proc)
        while not stop_dismisser:
            user32.EnumWindows(callback, 0)
            time.sleep(0.5)

    dismisser_thread = threading.Thread(target=dismiss_word_dialogs, daemon=True)
    dismisser_thread.start()

    try:
        progress("正在打开 PDF 文件（Word 原生解析）...", 15)
        # Word 自动识别 PDF 格式，无需指定 Format 参数
        doc = word.Documents.Open(str(Path(pdf_path).absolute()))

        progress("正在转换为 Word 格式（保留完整排版）...", 50)
        # FileFormat=16 即 wdFormatDocumentDefault (.docx)
        doc.SaveAs2(str(Path(docx_path).absolute()), 16)
        progress("正在保存...", 90)
    finally:
        stop_dismisser = True
        dismisser_thread.join(timeout=1)
        if doc is not None:
            try:
                doc.Close(False)
            except Exception:
                logging.exception("Could not close Word PDF document")
        try:
            word.Quit()
        except Exception:
            logging.exception("Could not quit Word after PDF conversion")
        finally:
            comtypes.CoUninitialize()
    progress("完成", 100)


def presentation_to_pdf_via_libreoffice(
    presentation_path: str, pdf_path: str, progress
) -> None:
    """Export PPT/PPTX to PDF through an isolated LibreOffice profile."""
    import shutil
    import tempfile

    lo = _get_lo_path()
    if not lo:
        raise RuntimeError(
            "LibreOffice 未安装或未找到。\n"
            "下载: https://www.libreoffice.org/download/"
        )
    progress("正在通过 LibreOffice 导出 PowerPoint...", 10)
    lo_dir = os.path.dirname(lo) or None

    with tempfile.TemporaryDirectory(prefix="pdf_word_converter_") as temp_root:
        temp_root_path = Path(temp_root)
        temp_output_dir = temp_root_path / "output"
        profile_dir = temp_root_path / "profile"
        temp_output_dir.mkdir()
        profile_dir.mkdir()
        profile_arg = f"-env:UserInstallation={profile_dir.as_uri()}"
        result = subprocess.run(
            [
                lo, profile_arg, "--headless", "--convert-to", "pdf",
                "--outdir", str(temp_output_dir),
                str(Path(presentation_path).absolute()),
            ],
            capture_output=True, text=True, timeout=180,
            cwd=lo_dir,
        )
        _raise_for_libreoffice_error(result, lo)
        generated = _find_libreoffice_output(
            temp_output_dir, Path(presentation_path).stem, ".pdf"
        )
        shutil.move(str(generated), str(pdf_path))
    progress("完成", 100)


def pdf_to_word_via_images(pdf_path: str, docx_path: str, progress, dpi: int = 300) -> None:
    """
    [高保真观感] 每页 PDF 渲染为高清图片嵌入 DOCX
    视觉上通常非常接近原 PDF；文字一般不可编辑、不可检索。
    """
    fitz = _load_pymupdf()
    from docx import Document
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Inches, Cm, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    progress("正在读取 PDF...", 5)
    pdf_doc = fitz.open(pdf_path)
    try:
        total_pages = len(pdf_doc)
        if total_pages == 0:
            raise RuntimeError("PDF 没有可转换的页面")

        word_doc = Document()
        word_doc.core_properties.title = Path(pdf_path).stem

        settings = word_doc.settings.element
        if settings.find(qn("w:doNotAutoCompressPictures")) is None:
            settings.append(OxmlElement("w:doNotAutoCompressPictures"))

        for i in range(total_pages):
            progress(
                f"渲染第 {i + 1}/{total_pages} 页为图片...",
                int(5 + (i / total_pages) * 80),
            )

            page = pdf_doc[i]
            render_dpi = calculate_adaptive_dpi(
                page.rect.width,
                page.rect.height,
                base_dpi=dpi,
                max_dimension_px=6000,
                max_pixels=30_000_000,
            )
            pix = page.get_pixmap(
                dpi=render_dpi,
                colorspace=fitz.csRGB,
                alpha=False,
            )
            img_bytes = pix.tobytes("png")

            page_w_inch = page.rect.width / 72
            page_h_inch = page.rect.height / 72
            page_scale = min(1.0, 22.0 / max(page_w_inch, page_h_inch))
            page_w_inch *= page_scale
            page_h_inch *= page_scale

            if i == 0:
                section = word_doc.sections[0]
                para = word_doc.add_paragraph()
            else:
                # The section-break paragraph belongs to the previous section;
                # the page image must be placed in a new paragraph after it.
                section = word_doc.add_section()
                para = word_doc.add_paragraph()

            section.page_width = Inches(page_w_inch)
            section.page_height = Inches(page_h_inch)
            section.top_margin = Cm(0)
            section.bottom_margin = Cm(0)
            section.left_margin = Cm(0)
            section.right_margin = Cm(0)
            section.header_distance = Cm(0)
            section.footer_distance = Cm(0)

            para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            pf = para.paragraph_format
            pf.space_before = Cm(0)
            pf.space_after = Cm(0)
            pf.line_spacing = 1.0
            pf.keep_together = True

            run = para.add_run()
            run.font.size = Pt(1)
            stream = io.BytesIO(img_bytes)
            run.add_picture(
                stream,
                width=Inches(page_w_inch),
                height=Inches(page_h_inch),
            )

        progress("正在保存 Word 文档...", 90)
        word_doc.save(docx_path)
    finally:
        pdf_doc.close()
    progress("完成", 100)


def calculate_adaptive_dpi(
    page_width_points: float,
    page_height_points: float,
    base_dpi: int = 200,
    max_dimension_px: int = 3200,
    max_pixels: int | None = None,
) -> int:
    """Return a high-quality DPI while bounding unusually large PDF pages."""
    width_points = float(page_width_points)
    height_points = float(page_height_points)
    if width_points <= 0 or height_points <= 0:
        raise ValueError("PDF 页面尺寸无效")
    if base_dpi <= 0 or max_dimension_px <= 0:
        raise ValueError("渲染分辨率限制必须大于 0")
    if max_pixels is not None and max_pixels <= 0:
        raise ValueError("最大像素数必须大于 0")
    longest_points = max(width_points, height_points)
    projected = longest_points / 72 * base_dpi
    dimension_dpi = float(base_dpi)
    if projected > max_dimension_px:
        dimension_dpi = base_dpi * max_dimension_px / projected

    pixel_dpi = float(base_dpi)
    if max_pixels is not None:
        projected_pixels = (
            width_points / 72 * base_dpi
        ) * (
            height_points / 72 * base_dpi
        )
        if projected_pixels > max_pixels:
            pixel_dpi = (
                base_dpi * (max_pixels / projected_pixels) ** 0.5
            )
    safe_dpi = min(float(base_dpi), dimension_dpi, pixel_dpi)
    if safe_dpi < 1:
        raise ValueError(
            "PDF 页面尺寸过大，即使按 1 DPI 渲染也会超过安全像素限制"
        )
    return max(1, int(safe_dpi))


def _presentation_canvas_size(page_width_points: float, page_height_points: float):
    """Fit the first PDF page ratio into a PowerPoint-compatible canvas."""
    from pptx.util import Inches

    width = float(page_width_points)
    height = float(page_height_points)
    if width <= 0 or height <= 0:
        raise ValueError("PDF 页面尺寸无效")

    long_edge = 13.333
    if width >= height:
        width_inches = long_edge
        height_inches = max(1.0, long_edge * height / width)
    else:
        height_inches = long_edge
        width_inches = max(1.0, long_edge * width / height)
    return Inches(width_inches), Inches(height_inches)


def pdf_to_pptx_via_images(pdf_path: str, pptx_path: str, progress) -> None:
    """Render each PDF page as a centered, uncropped image on one slide."""
    fitz = _load_pymupdf()
    from pptx import Presentation

    progress("正在读取 PDF...", 5)
    pdf_doc = fitz.open(pdf_path)
    try:
        total_pages = len(pdf_doc)
        if total_pages == 0:
            raise RuntimeError("PDF 没有可转换的页面")

        first_rect = pdf_doc[0].rect
        presentation = Presentation()
        presentation.slide_width, presentation.slide_height = _presentation_canvas_size(
            first_rect.width, first_rect.height
        )
        blank_layout = presentation.slide_layouts[6]

        for index, page in enumerate(pdf_doc):
            progress(
                f"正在生成第 {index + 1}/{total_pages} 张幻灯片...",
                int(5 + (index / total_pages) * 82),
            )
            dpi = calculate_adaptive_dpi(page.rect.width, page.rect.height)
            pixmap = page.get_pixmap(dpi=dpi, alpha=False)
            image_stream = io.BytesIO(pixmap.tobytes("png"))

            slide = presentation.slides.add_slide(blank_layout)
            page_ratio = page.rect.width / page.rect.height
            canvas_ratio = presentation.slide_width / presentation.slide_height
            if page_ratio >= canvas_ratio:
                picture_width = presentation.slide_width
                picture_height = int(picture_width / page_ratio)
            else:
                picture_height = presentation.slide_height
                picture_width = int(picture_height * page_ratio)
            left = int((presentation.slide_width - picture_width) / 2)
            top = int((presentation.slide_height - picture_height) / 2)
            slide.shapes.add_picture(
                image_stream, left, top, width=picture_width, height=picture_height
            )

        progress("正在保存 PowerPoint...", 92)
        presentation.save(pptx_path)
    finally:
        pdf_doc.close()
    progress("完成", 100)


def pdf_to_word_via_libreoffice(pdf_path: str, docx_path: str, progress) -> None:
    """通过 LibreOffice 无头模式转换。"""
    import shutil
    import tempfile

    lo = _get_lo_path()
    if not lo:
        raise RuntimeError(
            "LibreOffice 未安装或未找到。\n"
            "下载: https://www.libreoffice.org/download/"
        )
    progress("正在通过 LibreOffice 转换...", 10)
    lo_dir = os.path.dirname(lo) or None

    with tempfile.TemporaryDirectory(prefix="pdf_word_converter_") as temp_root:
        temp_root_path = Path(temp_root)
        temp_output_dir = temp_root_path / "output"
        profile_dir = temp_root_path / "profile"
        temp_output_dir.mkdir()
        profile_dir.mkdir()
        profile_arg = f"-env:UserInstallation={profile_dir.as_uri()}"
        result = subprocess.run(
            [
                lo, profile_arg, "--headless", "--convert-to", "docx",
                "--outdir", str(temp_output_dir), str(Path(pdf_path).absolute()),
            ],
            capture_output=True, text=True, timeout=300,
            cwd=lo_dir,
        )
        _raise_for_libreoffice_error(result, lo)
        generated = _find_libreoffice_output(
            temp_output_dir, Path(pdf_path).stem, ".docx"
        )
        shutil.move(str(generated), str(docx_path))
    progress("完成", 100)


def _raise_for_libreoffice_error(result, lo_path) -> None:
    stdout = getattr(result, "stdout", "") or ""
    stderr = getattr(result, "stderr", "") or ""
    combined_output = "\n".join(
        part.strip() for part in (stdout, stderr) if part.strip()
    )
    if result.returncode == 0 and "Error:" not in combined_output:
        return
    details = combined_output or f"LibreOffice 返回码 {result.returncode}，未提供详细信息"
    raise RuntimeError(
        f"LibreOffice 转换失败。\n"
        f"路径: {lo_path}\n"
        f"错误信息: {details}"
    )


def is_libreoffice_export_filter_error(error: Exception) -> bool:
    """Identify LibreOffice PDF import/export failures suitable for image fallback."""
    message = str(error).lower()
    indicators = (
        "export filter",
        "no export filter",
        "sfxbasemodel::impl_store",
        "未生成预期的输出文件",
        "source format: pdf",
    )
    return any(indicator in message for indicator in indicators)


def _find_libreoffice_output(temp_dir, source_stem, suffix) -> Path:
    expected = Path(temp_dir) / f"{source_stem}{suffix}"
    if expected.is_file():
        return expected
    candidates = list(Path(temp_dir).glob(f"*{suffix}"))
    if len(candidates) == 1:
        return candidates[0]
    raise RuntimeError("LibreOffice 未生成预期的输出文件")


# ═══════════════════════════════════════════════════════════
#  Word → PDF 转换方法
# ═══════════════════════════════════════════════════════════

def docx_to_pdf_via_word(docx_path: str, pdf_path: str, progress) -> None:
    """通过 Microsoft Word COM 导出 PDF"""
    import comtypes
    import comtypes.client

    progress("正在启动 Microsoft Word...", 5)
    comtypes.CoInitialize()
    try:
        word = comtypes.client.CreateObject("Word.Application", dynamic=True)
    except Exception:
        comtypes.CoUninitialize()
        raise
    word.Visible = False
    word.DisplayAlerts = 0  # wdAlertsNone
    doc = None

    try:
        progress("正在打开 Word 文档...", 15)
        doc = word.Documents.Open(str(Path(docx_path).absolute()), False, True)

        progress("正在导出 PDF...", 50)
        # FileFormat=17 即 wdFormatPDF
        doc.SaveAs2(str(Path(pdf_path).absolute()), 17)
    finally:
        if doc is not None:
            try:
                doc.Close(False)
            except Exception:
                logging.exception("Could not close Word document")
        try:
            word.Quit()
        except Exception:
            logging.exception("Could not quit Word after PDF export")
        finally:
            comtypes.CoUninitialize()
    progress("完成", 100)


def presentation_to_pdf_via_powerpoint(
    presentation_path: str, pdf_path: str, progress
) -> None:
    """Export PPT/PPTX to PDF through Microsoft PowerPoint COM."""
    import comtypes
    import comtypes.client

    progress("正在启动 Microsoft PowerPoint...", 5)
    comtypes.CoInitialize()
    powerpoint = None
    presentation = None
    try:
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application", dynamic=True)
        powerpoint.DisplayAlerts = 1  # ppAlertsNone
        try:
            powerpoint.AutomationSecurity = 3  # msoAutomationSecurityForceDisable
        except Exception:
            logging.debug("PowerPoint AutomationSecurity is unavailable", exc_info=True)

        progress("正在打开 PowerPoint 文件...", 15)
        presentation = powerpoint.Presentations.Open(
            str(Path(presentation_path).absolute()), True, False, False
        )
        progress("正在导出 PDF...", 50)
        presentation.SaveAs(str(Path(pdf_path).absolute()), 32)  # ppSaveAsPDF
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                logging.exception("Could not close PowerPoint presentation")
        if powerpoint is not None:
            try:
                powerpoint.Quit()
            except Exception:
                logging.exception("Could not quit PowerPoint")
        comtypes.CoUninitialize()
    progress("完成", 100)


def docx_to_pdf_via_libreoffice(docx_path: str, pdf_path: str, progress) -> None:
    """通过 LibreOffice 导出 PDF。"""
    import shutil
    import tempfile

    lo = _get_lo_path()
    if not lo:
        raise RuntimeError(
            "LibreOffice 未安装或未找到。\n"
            "下载: https://www.libreoffice.org/download/"
        )
    progress("正在通过 LibreOffice 导出...", 10)
    lo_dir = os.path.dirname(lo) or None

    with tempfile.TemporaryDirectory(prefix="pdf_word_converter_") as temp_root:
        temp_root_path = Path(temp_root)
        temp_output_dir = temp_root_path / "output"
        profile_dir = temp_root_path / "profile"
        temp_output_dir.mkdir()
        profile_dir.mkdir()
        profile_arg = f"-env:UserInstallation={profile_dir.as_uri()}"
        result = subprocess.run(
            [
                lo, profile_arg, "--headless", "--convert-to", "pdf",
                "--outdir", str(temp_output_dir), str(Path(docx_path).absolute()),
            ],
            capture_output=True, text=True, timeout=120,
            cwd=lo_dir,
        )
        _raise_for_libreoffice_error(result, lo)
        generated = _find_libreoffice_output(
            temp_output_dir, Path(docx_path).stem, ".pdf"
        )
        shutil.move(str(generated), str(pdf_path))
    progress("完成", 100)


# Kept as a compatibility shim for callers of earlier releases. The old
# implementation rewrote table widths, typography and punctuation using
# document-specific guesses, which can damage an otherwise usable conversion.
# v0.5.1 preserves the conversion engine's DOCX byte-for-byte instead.

def fix_converted_docx(docx_path: str, progress=None) -> None:
    """Deprecated no-op that preserves the engine-produced DOCX unchanged."""
    if not Path(docx_path).is_file():
        raise FileNotFoundError(docx_path)
    if progress:
        progress("已保留转换引擎原始排版", 95)


# ═══════════════════════════════════════════════════════════
#  GUI 界面
# ═══════════════════════════════════════════════════════════

from batch_logic import BatchResult, deduplicate_paths, run_conversion_batch


class ConverterApp:
    TARGET_DISPLAY = {
        PDF_TARGET_WORD: "Word (.docx)",
        PDF_TARGET_POWERPOINT: "PowerPoint (.pptx)",
        "pdf": "PDF (.pdf)",
    }
    TARGET_KIND = {label: kind for kind, label in TARGET_DISPLAY.items()}
    STATUS_TEXT = {
        "pending": "等待",
        "running": "转换中",
        "success": "成功",
        "failed": "失败",
        "cancelled": "已取消",
    }

    def __init__(self, root, platform_services=None):
        self.root = root
        self.platform_services = platform_services or PLATFORM_SERVICES
        self.is_macos = self.platform_services.platform == "darwin"
        self._mac_word_backend = None
        self._mac_powerpoint_backend = None
        if self.is_macos:
            self._mac_word_backend = MacOSWordBackend(
                installed_checker=lambda: self.platform_services.native_app_installed("word")
            )
            self._mac_powerpoint_backend = MacOSPowerPointBackend(
                installed_checker=lambda: self.platform_services.native_app_installed("powerpoint")
            )

        self.root.title(f"PDF ↔ Word/PPT 批量转换工具 v{APP_VERSION}")
        self.root.geometry("960x820")
        self.root.minsize(840, 720)
        self.root.resizable(True, True)
        self.root.configure(bg="#f5f5f5")
        self.root.protocol("WM_DELETE_WINDOW", self._on_window_close)
        self.ui_font_family = tkfont.nametofont("TkDefaultFont").actual("family")
        self.mono_font_family = tkfont.nametofont("TkFixedFont").actual("family")

        self.input_paths = []
        self.source_kind = None
        self.output_paths = []
        self._tree_paths = {}
        self._avail_methods = {
            WORD_NATIVE: None,
            POWERPOINT_NATIVE: None,
            "word_com": None,
            "powerpoint_com": None,
            "images": None,
            "pptx": None,
            "libreoffice": None,
        }
        self._engine_statuses = {
            key: EngineStatus(EngineState.CHECKING)
            for key in (WORD_NATIVE, POWERPOINT_NATIVE, "images", "pptx", "libreoffice")
        }
        self._method_widgets = {}
        self._engine_detection_complete = False
        self._engine_generation = 0
        self._install_detection_generation = None
        self._pending_engine_checks = set()
        self._engine_elapsed_ms = {}
        self._word_registered = False
        self._powerpoint_registered = False
        self._libreoffice_detected = False
        self._is_converting = False
        self._is_installing = False
        self._winget_path = None
        self._cancel_event = threading.Event()
        self._cancel_wait_started_at = None
        self._cancel_wait_after_id = None
        self._close_after_batch = False
        self._setup_ui()
        self._detect_engines()

    def _setup_ui(self):
        title = tk.Label(
            self.root, text="PDF ↔ Word/PPT 批量转换工具",
            font=(self.ui_font_family, 18, "bold"), bg="#f5f5f5", fg="#222",
        )
        title.pack(pady=(10, 2))

        tk.Label(
            self.root, text="本地串行转换 · 支持批量队列 · 文件不会上传",
            font=(self.ui_font_family, 10), bg="#f5f5f5", fg="#666",
        ).pack(pady=(0, 6))

        file_frame = tk.LabelFrame(
            self.root, text=" 转换队列 ", font=(self.ui_font_family, 10),
            bg="#f5f5f5", fg="#333", padx=10, pady=6,
        )
        file_frame.pack(padx=24, fill="both")

        button_row = tk.Frame(file_frame, bg="#f5f5f5")
        button_row.pack(fill="x", pady=(0, 7))

        self.select_pdf_btn = tk.Button(
            button_row, text="添加 PDF", font=(self.ui_font_family, 10),
            width=12, command=lambda: self.select_files("pdf"),
            bg="#c0392b", fg="white", activebackground="#a93226", cursor="hand2",
        )
        self.select_pdf_btn.pack(side="left", padx=(0, 7))

        self.select_docx_btn = tk.Button(
            button_row, text="添加 Word", font=(self.ui_font_family, 10),
            width=12, command=lambda: self.select_files("word"),
            bg="#2471a3", fg="white", activebackground="#1f618d", cursor="hand2",
        )
        self.select_docx_btn.pack(side="left", padx=(0, 7))

        self.select_ppt_btn = tk.Button(
            button_row, text="添加 PowerPoint", font=(self.ui_font_family, 10),
            width=15, command=lambda: self.select_files("powerpoint"),
            bg="#d35400", fg="white", activebackground="#ba4a00", cursor="hand2",
        )
        self.select_ppt_btn.pack(side="left", padx=(0, 7))

        self.remove_btn = tk.Button(
            button_row, text="移除选中", font=(self.ui_font_family, 9),
            width=11, command=self.remove_selected,
        )
        self.remove_btn.pack(side="left", padx=(8, 5))

        self.clear_btn = tk.Button(
            button_row, text="清空", font=(self.ui_font_family, 9),
            width=8, command=self.clear_queue,
        )
        self.clear_btn.pack(side="left")

        self.queue_summary_var = tk.StringVar(value="尚未添加文件")
        tk.Label(
            button_row, textvariable=self.queue_summary_var,
            font=(self.ui_font_family, 9), bg="#f5f5f5", fg="#555",
        ).pack(side="right")

        target_row = tk.Frame(file_frame, bg="#f5f5f5")
        target_row.pack(fill="x", pady=(0, 7))
        tk.Label(
            target_row, text="PDF 转换目标：", font=(self.ui_font_family, 9),
            bg="#f5f5f5", fg="#333",
        ).pack(side="left")
        self.target_display_var = tk.StringVar(value=self.TARGET_DISPLAY[PDF_TARGET_WORD])
        self.target_zones_frame = tk.Frame(target_row, bg="#f5f5f5")
        self.target_zones_frame.pack(side="left", fill="x", expand=True, padx=(6, 0))
        self._target_zones = {}
        self._create_pdf_target_zone(
            PDF_TARGET_WORD, "PDF → Word (.docx)\n拖放 PDF 到这里",
        )
        self._create_pdf_target_zone(
            PDF_TARGET_POWERPOINT, "PDF → PowerPoint (.pptx)\n拖放 PDF 到这里",
        )
        self._refresh_target_zones()

        self.tree_frame = tk.Frame(file_frame, bg="#f5f5f5", highlightthickness=1)
        self.tree_frame.pack(fill="both", expand=True)
        columns = ("name", "folder", "status", "output")
        self.file_tree = ttk.Treeview(
            self.tree_frame, columns=columns, show="headings", height=4,
            selectmode="extended",
        )
        self.file_tree.heading("name", text="文件名")
        self.file_tree.heading("folder", text="来源目录")
        self.file_tree.heading("status", text="状态")
        self.file_tree.heading("output", text="输出文件")
        self.file_tree.column("name", width=180, minwidth=120, stretch=False)
        self.file_tree.column("folder", width=250, minwidth=140, stretch=True)
        self.file_tree.column("status", width=75, minwidth=70, stretch=False, anchor="center")
        self.file_tree.column("output", width=260, minwidth=160, stretch=True)
        self.file_tree.tag_configure("success", foreground="#18794e")
        self.file_tree.tag_configure("failed", foreground="#b42318")
        self.file_tree.tag_configure("cancelled", foreground="#777")
        tree_scroll = ttk.Scrollbar(self.tree_frame, command=self.file_tree.yview)
        self.file_tree.configure(yscrollcommand=tree_scroll.set)
        tree_scroll.pack(side="right", fill="y")
        self.file_tree.pack(side="left", fill="both", expand=True)
        self.file_tree.bind("<<TreeviewSelect>>", lambda _event: self._update_action_states())
        self.empty_drop_hint = tk.Label(
            self.tree_frame,
            text="拖放 PDF、Word 或 PowerPoint 文件到这里\nPDF 使用上方目标；Word/PPT 自动转换为 PDF",
            font=(self.ui_font_family, 10), bg="white", fg="#666",
            justify="center", padx=18, pady=10,
        )
        self.empty_drop_hint.place(relx=0.5, rely=0.56, anchor="center")
        self._setup_drag_and_drop()

        output_frame = tk.LabelFrame(
            self.root, text=" 输出位置 ", font=(self.ui_font_family, 10),
            bg="#f5f5f5", fg="#333", padx=10, pady=8,
        )
        output_frame.pack(padx=24, pady=(6, 0), fill="x")

        self.output_mode_var = tk.StringVar(value="source")
        self.source_output_radio = tk.Radiobutton(
            output_frame, text="各源文件所在目录", variable=self.output_mode_var,
            value="source", command=self._toggle_output_mode,
            font=(self.ui_font_family, 9), bg="#f5f5f5",
        )
        self.source_output_radio.pack(side="left", padx=(0, 12))
        self.custom_output_radio = tk.Radiobutton(
            output_frame, text="统一输出目录", variable=self.output_mode_var,
            value="custom", command=self._toggle_output_mode,
            font=(self.ui_font_family, 9), bg="#f5f5f5",
        )
        self.custom_output_radio.pack(side="left")

        self.output_dir_var = tk.StringVar()
        self.output_entry = tk.Entry(
            output_frame, textvariable=self.output_dir_var,
            font=(self.mono_font_family, 9), state="readonly", readonlybackground="white",
        )
        self.output_entry.pack(side="left", padx=8, fill="x", expand=True)
        self.browse_output_btn = tk.Button(
            output_frame, text="选择...", font=(self.ui_font_family, 9),
            width=9, command=self.choose_output_dir,
        )
        self.browse_output_btn.pack(side="right")

        self.method_frame = tk.LabelFrame(
            self.root, text=" 转换方式 ", font=(self.ui_font_family, 10),
            bg="#f5f5f5", fg="#333", padx=10, pady=5,
        )
        self.method_frame.pack(padx=24, pady=(6, 0), fill="x")

        self.method_var = tk.StringVar(value="images" if self.is_macos else WORD_NATIVE)
        self._method_labels = {}
        self.word_method_container = tk.Frame(self.method_frame, bg="#f5f5f5")
        methods_desc = []
        if not self.is_macos:
            methods_desc.append((WORD_NATIVE, "可编辑优先：Microsoft Word 原生转换（复杂表格可能重排）"))
        methods_desc.extend([
            ("images", "仅扫描件/纯图形：整页高清图片（文字和表格不可编辑）"),
            ("libreoffice", "兼容模式：LibreOffice（可编辑性与版式有限）"),
        ])
        for value, description in methods_desc:
            row = tk.Frame(self.word_method_container, bg="#f5f5f5")
            row.pack(anchor="w", pady=1, fill="x")
            radio = tk.Radiobutton(
                row, text=description, variable=self.method_var, value=value,
                font=(self.ui_font_family, 9), bg="#f5f5f5", anchor="w",
                command=lambda selected=value: self._on_method_selected(selected),
            )
            radio.pack(side="left")
            status_label = tk.Label(
                row, text="检测中...", font=(self.ui_font_family, 8),
                bg="#f5f5f5", fg="#777", width=10, anchor="e",
            )
            status_label.pack(side="right")
            self._method_widgets[value] = radio
            self._method_labels[value] = status_label

        self.ppt_method_container = tk.Frame(self.method_frame, bg="#f5f5f5")
        tk.Label(
            self.ppt_method_container,
            text="内置自适应高清图片模式（推荐，整页图片不可编辑）",
            font=(self.ui_font_family, 9), bg="#f5f5f5", anchor="w",
        ).pack(side="left")
        self.pptx_status_label = tk.Label(
            self.ppt_method_container, text="检测中...", font=(self.ui_font_family, 8),
            bg="#f5f5f5", fg="#777", width=10, anchor="e",
        )
        self.pptx_status_label.pack(side="right")

        self.auto_method_container = tk.Frame(self.method_frame, bg="#f5f5f5")
        self.auto_method_var = tk.StringVar(value="自动选择可用引擎")
        tk.Label(
            self.auto_method_container, textvariable=self.auto_method_var,
            font=(self.ui_font_family, 9), bg="#f5f5f5", anchor="w",
        ).pack(anchor="w")

        self.engine_status_var = tk.StringVar(value="环境检测中...")
        tk.Label(
            self.method_frame, textvariable=self.engine_status_var,
            font=(self.ui_font_family, 8), bg="#f5f5f5", fg="#666", anchor="w",
        ).pack(side="bottom", fill="x", pady=(5, 0))
        self._update_method_panel()

        action_frame = tk.Frame(self.root, bg="#f5f5f5")
        action_frame.pack(pady=(6, 4))
        self.convert_btn = tk.Button(
            action_frame, text="开始批量转换", font=(self.ui_font_family, 10, "bold"),
            width=15, command=self.start_conversion,
            bg="#1e8449", fg="white", activebackground="#196f3d", cursor="hand2",
        )
        self.convert_btn.pack(side="left", padx=5)
        self.cancel_btn = tk.Button(
            action_frame, text="取消", font=(self.ui_font_family, 10),
            width=10, command=self.cancel_conversion,
            bg="#a04000", fg="white", activebackground="#873600", cursor="hand2",
        )
        self.cancel_btn.pack(side="left", padx=5)
        self.open_btn = tk.Button(
            action_frame, text="打开输出目录", font=(self.ui_font_family, 10),
            width=14, command=self.open_output_folder,
            bg="#555", fg="white", activebackground="#333", cursor="hand2",
        )
        self.open_btn.pack(side="left", padx=5)

        progress_frame = tk.Frame(self.root, bg="#f5f5f5")
        progress_frame.pack(padx=24, fill="x")
        self.progress = ttk.Progressbar(progress_frame, mode="determinate", maximum=100)
        self.progress.pack(fill="x")
        self.progress_text_var = tk.StringVar(value="就绪")
        tk.Label(
            progress_frame, textvariable=self.progress_text_var,
            font=(self.ui_font_family, 8), bg="#f5f5f5", fg="#666", anchor="w",
        ).pack(fill="x", pady=(2, 0))

        log_frame = tk.LabelFrame(
            self.root, text=" 转换日志 ", font=(self.ui_font_family, 10),
            bg="#f5f5f5", fg="#333", padx=7, pady=4,
        )
        log_frame.pack(padx=24, pady=(4, 8), fill="both", expand=True)
        self.log_text = tk.Text(
            log_frame, height=7, font=(self.mono_font_family, 9),
            bg="#1e1e1e", fg="#d4d4d4", wrap="word", state="disabled",
        )
        log_scroll = ttk.Scrollbar(log_frame, command=self.log_text.yview)
        self.log_text.configure(yscrollcommand=log_scroll.set)
        log_scroll.pack(side="right", fill="y")
        self.log_text.pack(side="left", fill="both", expand=True)
        self._update_action_states()

    def _create_pdf_target_zone(self, target_kind, text):
        zone = tk.Label(
            self.target_zones_frame,
            text=text,
            font=(self.ui_font_family, 9, "bold"),
            bg="white",
            fg="#333",
            borderwidth=1,
            relief="solid",
            padx=10,
            pady=6,
            cursor="hand2",
        )
        zone.pack(side="left", fill="x", expand=True, padx=(0, 6))
        zone.bind(
            "<Button-1>",
            lambda _event, kind=target_kind: self._select_pdf_target(kind),
        )
        self._target_zones[target_kind] = zone

    def _refresh_target_zones(self):
        if not hasattr(self, "_target_zones"):
            return
        selected = self._current_target_kind()
        editable = (
            not self._is_converting
            and not self._is_installing
            and self.source_kind in (None, "pdf")
        )
        for target_kind, zone in self._target_zones.items():
            is_selected = target_kind == selected
            zone.config(
                bg="#e8f5e9" if is_selected else ("white" if editable else "#ededed"),
                fg="#145a32" if is_selected else ("#333" if editable else "#888"),
                relief="sunken" if is_selected else "solid",
                cursor="hand2" if editable else "arrow",
            )

    def _confirm_pdf_target_change(self, target_kind):
        current_target = self._current_target_kind()
        if (
            self.source_kind != "pdf"
            or not self.input_paths
            or current_target == target_kind
        ):
            return True
        return messagebox.askyesno(
            "切换整批转换目标",
            f"当前队列中的 {len(self.input_paths)} 个 PDF 将全部从 "
            f"{self.TARGET_DISPLAY[current_target]} 改为 "
            f"{self.TARGET_DISPLAY[target_kind]}。\n\n是否继续？",
        )

    def _select_pdf_target(self, target_kind):
        if self._is_converting or self._is_installing:
            return
        if self.source_kind not in (None, "pdf"):
            return
        if not self._confirm_pdf_target_change(target_kind):
            return
        self._set_target_kind(target_kind)
        self._update_method_panel()
        self._refresh_queue()

    def _setup_drag_and_drop(self):
        dnd_ready = (
            DND_FILES is not None
            and getattr(self.root, "_pdf_converter_dnd_available", True)
        )
        if not dnd_ready:
            self.empty_drop_hint.config(
                text="拖放组件不可用，请使用上方按钮添加文件",
                fg="#666",
            )
            return

        try:
            for widget in (self.file_tree, self.empty_drop_hint):
                widget.drop_target_register(DND_FILES)
                widget.dnd_bind(
                    "<<DropEnter>>",
                    lambda event: self._on_drop_hover(event, None, True),
                )
                widget.dnd_bind(
                    "<<DropLeave>>",
                    lambda event: self._on_drop_hover(event, None, False),
                )
                widget.dnd_bind(
                    "<<Drop>>",
                    lambda event: self._on_files_dropped(event, None),
                )

            for target_kind, widget in self._target_zones.items():
                widget.drop_target_register(DND_FILES)
                widget.dnd_bind(
                    "<<DropEnter>>",
                    lambda event, kind=target_kind: self._on_drop_hover(
                        event, kind, True,
                    ),
                )
                widget.dnd_bind(
                    "<<DropLeave>>",
                    lambda event, kind=target_kind: self._on_drop_hover(
                        event, kind, False,
                    ),
                )
                widget.dnd_bind(
                    "<<Drop>>",
                    lambda event, kind=target_kind: self._on_files_dropped(
                        event, kind,
                    ),
                )
        except (AttributeError, tk.TclError):
            logging.exception("TkDND could not be initialized; falling back to file buttons")
            self.root._pdf_converter_dnd_available = False
            self.empty_drop_hint.config(
                text="拖放组件不可用，请使用上方按钮添加文件",
                fg="#666",
            )

    def _on_drop_hover(self, event, target_kind, active):
        if target_kind is None:
            color = "#1e8449" if active else "#d0d0d0"
            self.tree_frame.config(highlightbackground=color, highlightcolor=color)
        elif active:
            self._target_zones[target_kind].config(bg="#fff2cc")
        else:
            self._refresh_target_zones()
        if self._is_converting or self._is_installing:
            return REFUSE_DROP
        return COPY

    def _on_files_dropped(self, event, forced_pdf_target=None):
        self._on_drop_hover(event, forced_pdf_target, False)
        if self._is_converting or self._is_installing:
            return REFUSE_DROP

        try:
            dropped = self.root.tk.splitlist(event.data)
        except tk.TclError as exc:
            self.log(f"无法解析拖入的文件: {exc}")
            messagebox.showerror("拖放失败", "无法读取拖入的文件路径。")
            return REFUSE_DROP

        try:
            source_kind, accepted, ignored = classify_dropped_paths(dropped)
        except MixedSourceKindsError as exc:
            self.log(str(exc))
            messagebox.showwarning(
                "请分批拖入",
                f"{exc}\n\n每个批次只能处理一种源格式，请分开拖入。",
            )
            return REFUSE_DROP

        if not accepted:
            messagebox.showwarning(
                "没有可添加的文件",
                "仅支持 PDF、DOCX、PPT 和 PPTX 文件；文件夹不会递归导入。",
            )
            return REFUSE_DROP

        if forced_pdf_target is not None and source_kind != "pdf":
            messagebox.showinfo(
                "此区域仅接收 PDF",
                "请将 Word 或 PowerPoint 文件拖到下方转换队列。",
            )
            return REFUSE_DROP

        if source_kind == "pdf":
            target_kind = forced_pdf_target or self._current_target_kind()
            if target_kind == "pdf":
                target_kind = PDF_TARGET_WORD
        else:
            target_kind = "pdf"

        added = self._add_paths(source_kind, accepted, target_kind)
        if not added:
            return REFUSE_DROP
        if added and ignored:
            messagebox.showwarning(
                "已忽略部分内容",
                f"已添加 {added} 个文件，另有 {len(ignored)} 个不支持的文件或文件夹被忽略。",
            )
        return COPY

    def _detect_engines(self, force=False):
        if force:
            reset_engine_cache()

        if self.is_macos:
            self._word_registered = self.platform_services.native_app_installed("word")
            self._powerpoint_registered = self.platform_services.native_app_installed(
                "powerpoint"
            )
            word_probe = self._mac_word_backend.probe
            powerpoint_probe = self._mac_powerpoint_backend.probe
        else:
            self._word_registered = _com_application_registered("Word.Application")
            self._powerpoint_registered = _com_application_registered(
                "PowerPoint.Application"
            )
            word_probe = word_com_available
            powerpoint_probe = powerpoint_com_available
        self._libreoffice_detected = bool(_find_libreoffice())
        self._engine_generation += 1
        generation = self._engine_generation
        if force and self._is_installing:
            self._install_detection_generation = generation
        self._engine_detection_complete = False
        self._engine_started_at = time.perf_counter()
        probes = {
            WORD_NATIVE: word_probe,
            POWERPOINT_NATIVE: powerpoint_probe,
            "libreoffice": libreoffice_available,
            "images": pymupdf_available,
            "pptx": pptx_component_available,
        }
        if not self.is_macos:
            probes["winget"] = find_winget
        self._pending_engine_checks = set(probes)
        self._engine_elapsed_ms = {}
        for key in self._avail_methods:
            self._avail_methods[key] = None
        for key in (WORD_NATIVE, POWERPOINT_NATIVE, "libreoffice", "images", "pptx"):
            self._engine_statuses[key] = EngineStatus(EngineState.CHECKING)
        self._winget_path = None
        self._update_engine_indicators()
        self._refresh_engine_status_text()
        self._update_action_states()

        for key, probe in probes.items():
            threading.Thread(
                target=self._run_engine_probe,
                args=(generation, key, probe),
                daemon=True,
            ).start()
            self.root.after(
                20000,
                self._expire_engine_probe,
                generation,
                key,
            )

    def _expire_engine_probe(self, generation, key):
        if generation != self._engine_generation:
            return
        if key not in self._pending_engine_checks:
            return
        self._apply_engine_probe_result(
            generation,
            key,
            EngineStatus(
                EngineState.TIMEOUT,
                "检测超时，后台检查已停止等待",
            ) if key != "winget" else False,
            20000,
        )

    def _run_engine_probe(self, generation, key, probe):
        started_at = time.perf_counter()
        error = None
        try:
            result = probe()
        except Exception as exc:
            logging.exception("Environment probe failed: %s", key)
            result = False
            error = str(exc)
        elapsed_ms = round((time.perf_counter() - started_at) * 1000)
        self.root.after(
            0,
            self._apply_engine_probe_result,
            generation,
            key,
            result,
            elapsed_ms,
            error,
        )

    def _apply_engine_probe_result(
        self, generation, key, result, elapsed_ms, error=None,
    ):
        if generation != self._engine_generation:
            return
        if key not in self._pending_engine_checks:
            return

        self._pending_engine_checks.discard(key)
        self._engine_elapsed_ms[key] = elapsed_ms
        if key == "winget":
            self._winget_path = result or None
            available = bool(self._winget_path)
        else:
            status = _coerce_engine_status(result, error)
            self._set_engine_status(key, status)
            available = self._method_availability(key)

        display_names = {
            WORD_NATIVE: "Microsoft Word",
            POWERPOINT_NATIVE: "Microsoft PowerPoint",
            "word_com": "Microsoft Word",
            "powerpoint_com": "Microsoft PowerPoint",
            "libreoffice": "LibreOffice",
            "images": "PyMuPDF 图片模式",
            "pptx": "内置 PPTX 组件",
            "winget": "winget",
        }
        status_detail = ""
        if key != "winget":
            status_detail = self._engine_statuses[_canonical_engine_key(key)].detail
        detail_text = status_detail or (str(error) if error else "")
        detail = f"；{detail_text}" if detail_text else ""
        self.log(
            f"{display_names[key]}: {self._availability_text(key, available)}"
            f"（检测 {elapsed_ms} ms）{detail}"
        )

        current_method = self.method_var.get()
        if (
            not self._is_installing
            and self._method_availability(current_method) is False
        ):
            self._select_best_pdf_method()

        if not self._pending_engine_checks:
            self._engine_detection_complete = True
            if self._install_detection_generation == generation:
                self._is_installing = False
                self._install_detection_generation = None
            total_ms = round((time.perf_counter() - self._engine_started_at) * 1000)
            self.log(f"环境检测完成，共 {total_ms} ms")

        self._update_engine_indicators()
        self._refresh_engine_status_text()
        self._update_method_panel()
        self._update_action_states()

    def _availability_text(self, key, available):
        canonical = _canonical_engine_key(key)
        status = getattr(self, "_engine_statuses", {}).get(canonical)
        if status is not None:
            labels = {
                EngineState.CHECKING: "检测中",
                EngineState.AVAILABLE: "可用",
                EngineState.MISSING: "不可用",
                EngineState.UNVERIFIED: "已安装/使用时验证",
                EngineState.PERMISSION_DENIED: "需要自动化权限",
                EngineState.LAUNCH_FAILED: "已安装/启动失败",
                EngineState.TIMEOUT: "检测超时",
                EngineState.UNSUPPORTED: "此平台不支持",
            }
            return labels[status.state]
        if available is None:
            return "检测中"
        if available:
            return "可用"
        detected = {
            WORD_NATIVE: getattr(self, "_word_registered", False),
            POWERPOINT_NATIVE: getattr(self, "_powerpoint_registered", False),
            "word_com": getattr(self, "_word_registered", False),
            "powerpoint_com": getattr(self, "_powerpoint_registered", False),
            "libreoffice": getattr(self, "_libreoffice_detected", False),
        }.get(key, False)
        return "已安装/启动失败" if detected else "不可用"

    def _set_engine_status(self, key, status):
        canonical = _canonical_engine_key(key)
        if not hasattr(self, "_engine_statuses"):
            self._engine_statuses = {}
        self._engine_statuses[canonical] = status
        availability = None if status.state is EngineState.CHECKING else _status_selectable(status)
        self._avail_methods[canonical] = availability
        if canonical == WORD_NATIVE:
            self._avail_methods["word_com"] = availability
        elif canonical == POWERPOINT_NATIVE:
            self._avail_methods["powerpoint_com"] = availability

    def _method_availability(self, key):
        canonical = _canonical_engine_key(key)
        if canonical in self._avail_methods:
            return self._avail_methods.get(canonical)
        if key in self._avail_methods:
            return self._avail_methods.get(key)
        legacy = {
            WORD_NATIVE: "word_com",
            POWERPOINT_NATIVE: "powerpoint_com",
        }.get(canonical)
        return self._avail_methods.get(legacy) if legacy else None

    def _update_engine_indicators(self):
        for value, label in self._method_labels.items():
            available = self._method_availability(value)
            status_text = self._availability_text(value, available)
            if available is None:
                label.config(text="检测中...", fg="#777")
            elif status_text == "已安装/使用时验证":
                label.config(text="使用时验证", fg="#a04000")
            elif status_text == "需要自动化权限":
                label.config(text="需要权限", fg="#b42318")
            elif available:
                label.config(text="✓ 可用", fg="#18794e")
            elif status_text == "已安装/启动失败":
                label.config(text="启动失败", fg="#a04000")
            elif status_text == "检测超时":
                label.config(text="检测超时", fg="#a04000")
            else:
                label.config(text="✗ 不可用", fg="#b42318")

        image_state = self._avail_methods.get("images")
        pptx_state = self._avail_methods.get("pptx")
        if image_state is None or pptx_state is None:
            self.pptx_status_label.config(text="检测中...", fg="#777")
        else:
            available = image_state and pptx_state
            self.pptx_status_label.config(
                text="✓ 可用" if available else "✗ 不可用",
                fg="#18794e" if available else "#b42318",
            )

    def _refresh_engine_status_text(self):
        native_line = " | ".join((
            f"Word: {self._availability_text(WORD_NATIVE, self._method_availability(WORD_NATIVE))}",
            f"PowerPoint: {self._availability_text(POWERPOINT_NATIVE, self._method_availability(POWERPOINT_NATIVE))}",
            f"LibreOffice: {self._availability_text('libreoffice', self._method_availability('libreoffice'))}",
        ))
        component_parts = [
            f"PyMuPDF: {self._availability_text('images', self._method_availability('images'))}",
            f"PPTX组件: {self._availability_text('pptx', self._method_availability('pptx'))}",
        ]
        if not getattr(self, "is_macos", False):
            winget_status = (
                "检测中"
                if "winget" in self._pending_engine_checks
                else ("可用" if self._winget_path else "不可用")
            )
            component_parts.append(f"winget: {winget_status}")
        self.engine_status_var.set(
            native_line + "\n" + " | ".join(component_parts)
        )

    def _update_engine_status(
        self, have_word, have_powerpoint, have_libreoffice,
        have_pymupdf, have_pptx, winget_path,
    ):
        """Apply a complete result set; retained for logic tests and integrations."""
        self._avail_methods = {
            WORD_NATIVE: _availability_from_result(have_word),
            POWERPOINT_NATIVE: _availability_from_result(have_powerpoint),
            "word_com": _availability_from_result(have_word),
            "powerpoint_com": _availability_from_result(have_powerpoint),
            "images": _availability_from_result(have_pymupdf),
            "pptx": _availability_from_result(have_pptx),
            "libreoffice": _availability_from_result(have_libreoffice),
        }
        self._engine_statuses = {
            WORD_NATIVE: _coerce_engine_status(have_word),
            POWERPOINT_NATIVE: _coerce_engine_status(have_powerpoint),
            "libreoffice": _coerce_engine_status(have_libreoffice),
            "images": _coerce_engine_status(have_pymupdf),
            "pptx": _coerce_engine_status(have_pptx),
        }
        self._winget_path = winget_path
        self._pending_engine_checks = set()
        self._engine_detection_complete = True
        self._select_best_pdf_method()
        self._update_engine_indicators()
        self._refresh_engine_status_text()
        self._update_method_panel()
        self._update_action_states()

    def _current_target_kind(self):
        return self.TARGET_KIND.get(self.target_display_var.get(), PDF_TARGET_WORD)

    def _set_target_kind(self, target_kind):
        self.target_display_var.set(self.TARGET_DISPLAY[target_kind])
        if hasattr(self, "_target_zones"):
            self._refresh_target_zones()

    def _current_spec(self):
        if not self.source_kind:
            return None
        pdf_target = self._current_target_kind()
        return resolve_conversion_spec(self.source_kind, pdf_target)

    def _conversion_detection_ready(self):
        spec = self._current_spec()
        if spec is None:
            return False
        if spec.key == "pdf_to_word":
            selected_state = self._method_availability(self.method_var.get())
            if self.method_var.get() == "images":
                if selected_state is None:
                    return False
                libreoffice_state = self._method_availability("libreoffice")
                if self.is_macos:
                    return libreoffice_state is not None
                word_state = self._method_availability(WORD_NATIVE)
                return word_state is True or (
                    word_state is False and libreoffice_state is not None
                )
            return (
                selected_state is not None
                and self._method_availability("images") is not None
            )
        if spec.key == "pdf_to_powerpoint":
            return (
                self._method_availability("images") is not None
                and self._method_availability("pptx") is not None
            )
        if spec.key == "word_to_pdf":
            word_state = self._method_availability(WORD_NATIVE)
            libreoffice_state = self._method_availability("libreoffice")
            return word_state is True or (
                word_state is False and libreoffice_state is not None
            )
        if spec.key == "powerpoint_to_pdf":
            powerpoint_state = self._method_availability(POWERPOINT_NATIVE)
            libreoffice_state = self._method_availability("libreoffice")
            return powerpoint_state is True or (
                powerpoint_state is False and libreoffice_state is not None
            )
        return False

    def _update_method_panel(self):
        for container in (
            self.word_method_container,
            self.ppt_method_container,
            self.auto_method_container,
        ):
            container.pack_forget()

        source_kind = self.source_kind or "pdf"
        pdf_target = self._current_target_kind()
        if source_kind == "pdf" and pdf_target == "pdf":
            pdf_target = PDF_TARGET_WORD
        spec = resolve_conversion_spec(source_kind, pdf_target)
        self.method_frame.config(text=f" {spec.label.replace(' -> ', ' → ')} 转换方式 ")
        if spec.key == "pdf_to_word":
            self.word_method_container.pack(fill="x")
        elif spec.key == "pdf_to_powerpoint":
            self.ppt_method_container.pack(fill="x")
        else:
            native_kind = "AppleScript" if getattr(self, "is_macos", False) else "本机 Office"
            if spec.key == "word_to_pdf":
                text = f"自动选择：Microsoft Word {native_kind} 优先，LibreOffice 备用"
            else:
                text = f"自动选择：Microsoft PowerPoint {native_kind} 优先，LibreOffice 备用"
            self.auto_method_var.set(text)
            self.auto_method_container.pack(fill="x")

    def _on_target_selected(self, _event=None):
        if self.source_kind and self.source_kind != "pdf":
            self._set_target_kind("pdf")
        self._update_method_panel()
        self._refresh_queue()

    def _select_best_pdf_method(self):
        methods = ("images", "libreoffice") if getattr(self, "is_macos", False) else (
            WORD_NATIVE, "images", "libreoffice"
        )
        for method in methods:
            if self._method_availability(method):
                self.method_var.set(method)
                return

    def _on_method_selected(self, method):
        if method == WORD_NATIVE and not self._method_availability(WORD_NATIVE):
            accepted = self._prompt_word_install()
            if not accepted:
                self._select_best_pdf_method()
        elif method == "libreoffice" and not self._method_availability("libreoffice"):
            accepted = self._prompt_libreoffice_install()
            if not accepted:
                self._select_best_pdf_method()

    def _prompt_word_install(self):
        if self._method_availability(WORD_NATIVE) or self._is_installing:
            return False
        if getattr(self, "_word_registered", False):
            self.log("检测到 Microsoft Word 已安装，但原生自动化暂不可用")
            if getattr(self, "is_macos", False):
                guidance = (
                    "请前往“系统设置 → 隐私与安全性 → 自动化”，允许本工具控制 "
                    "Microsoft Word，然后重试；无需重新安装 Office。"
                )
            else:
                guidance = (
                    "请关闭 Word 的登录、许可或文档提示窗口后重新启动本工具，"
                    "无需重复安装 Office。"
                )
            messagebox.showwarning(
                "Microsoft Word 暂不可用",
                "已检测到 Microsoft Word，但自动化启动失败。\n\n" + guidance,
            )
            return False
        install_text = (
            "是否打开 Microsoft Office 官方下载页？"
            if getattr(self, "is_macos", False)
            else "是否使用 winget 安装 Microsoft Office？"
        )
        confirmed = messagebox.askyesno(
            "未检测到 Microsoft Word",
            install_text + "\n\n"
            "Office 需要有效许可证和 Microsoft 账号，安装过程可能要求管理员权限。"
            "程序不会自动提供许可证。",
        )
        if confirmed:
            self._start_install("office")
        else:
            self.log("用户已拒绝安装 Microsoft Office")
        return confirmed

    def _prompt_libreoffice_install(self):
        if self._method_availability("libreoffice") or self._is_installing:
            return False
        if getattr(self, "_libreoffice_detected", False):
            self.log("检测到 LibreOffice 已安装，但无头模式启动失败")
            messagebox.showwarning(
                "LibreOffice 暂不可用",
                "已检测到 LibreOffice 程序，但后台启动失败。\n\n"
                "请关闭正在运行的 LibreOffice 窗口后重新启动本工具，"
                "无需重复安装。",
            )
            return False
        install_text = (
            "是否打开 LibreOffice 官方下载页？"
            if getattr(self, "is_macos", False)
            else "是否使用 winget 安装官方 LibreOffice？"
        )
        confirmed = messagebox.askyesno(
            "需要 LibreOffice",
            "当前任务没有可用的 LibreOffice 引擎。" + install_text,
        )
        if confirmed:
            self._start_install("libreoffice")
        else:
            self.log("用户已拒绝安装 LibreOffice")
        return confirmed

    def _prompt_macos_conversion_engines(self, native_name, native_key, registered):
        status = getattr(self, "_engine_statuses", {}).get(native_key)
        if registered:
            detail = status.detail if status and status.detail else self._availability_text(
                native_key,
                self._method_availability(native_key),
            )
            if status and status.state is EngineState.PERMISSION_DENIED:
                guidance = "请先检查“系统设置 → 隐私与安全性 → 自动化”的授权。"
            elif status and status.state is EngineState.TIMEOUT:
                guidance = "请关闭 Office 中的弹窗或占用文件后重试。"
            else:
                guidance = "请先确认 Office 已完成登录和许可证激活。"
            use_libreoffice = messagebox.askyesno(
                f"{native_name} 暂不可用",
                f"已检测到 {native_name}，但原生转换暂不可用。\n\n"
                f"{guidance}\n{detail}\n\n"
                "是否打开 LibreOffice 官方下载页作为备用？",
            )
            if use_libreoffice:
                self._start_install("libreoffice")
            return use_libreoffice

        choice = messagebox.askyesnocancel(
            "缺少转换引擎",
            f"未检测到 {native_name} 或 LibreOffice。\n\n"
            f"选择“是”打开 Microsoft Office 官方下载页；\n"
            "选择“否”打开 LibreOffice 官方下载页；\n"
            "选择“取消”则不执行任何操作。\n\n"
            "程序不会自动安装软件，也不会提供 Office 许可证。",
        )
        if choice is None:
            self.log("用户取消打开转换引擎下载页")
            return False
        product = "office" if choice else "libreoffice"
        self._start_install(product)
        return True

    def _start_install(self, product):
        display_name = "Microsoft Office" if product == "office" else "LibreOffice"
        action = self.platform_services.installer_action(product)
        if action.kind is InstallerActionKind.OPEN_URL:
            self.log(f"正在打开 {display_name} 官方下载页")
            result = self.platform_services.perform_installer_action(action)
            if not result.succeeded:
                messagebox.showerror(
                    "无法打开下载页",
                    f"无法打开 {display_name} 官方下载页。\n\n{result.detail}",
                )
            return

        self._is_installing = True
        self._engine_detection_complete = False
        self._update_action_states()
        self.progress_text_var.set(f"正在安装 {display_name}...")
        self.log(f"开始安装 {display_name}")

        def install():
            result = self.platform_services.perform_installer_action(action)
            self.root.after(0, self._finish_install, product, display_name, result)

        threading.Thread(target=install, daemon=True).start()

    def _finish_install(self, product, display_name, result):
        if result.succeeded:
            self.log(f"{display_name} 安装命令已完成，正在重新检测引擎")
            messagebox.showinfo("安装完成", f"{display_name} 安装已完成，程序将立即重新检测。")
            self._detect_engines(force=True)
            return

        self._is_installing = False
        self._engine_detection_complete = True
        self._update_action_states()
        failure_detail = getattr(result, "message", "") or getattr(result, "detail", "")
        self.log(f"{display_name} 安装失败: {failure_detail}")
        messagebox.showwarning(
            "安装失败",
            f"{display_name} 安装未成功。将打开官方下载页。\n\n{failure_detail}",
        )
        open_official_download(product)

    def _ask_yes_no_from_worker(self, title, message):
        completed = threading.Event()
        answer = {"value": False}

        def ask():
            try:
                answer["value"] = messagebox.askyesno(title, message)
            finally:
                completed.set()

        self.root.after(0, ask)
        completed.wait()
        return answer["value"]

    def _cancel_requested(self):
        cancel_event = getattr(self, "_cancel_event", None)
        return bool(cancel_event and cancel_event.is_set())

    def _ask_yes_no_cancel_from_worker(self, title, message):
        if self._cancel_requested():
            return None
        completed = threading.Event()
        answer = {"value": None}

        def ask():
            if self._cancel_requested():
                completed.set()
                return
            try:
                answer["value"] = messagebox.askyesnocancel(title, message)
            finally:
                completed.set()

        self.root.after(0, ask)
        completed.wait()
        return None if self._cancel_requested() else answer["value"]

    def _ask_ok_cancel_from_worker(self, title, message):
        if self._cancel_requested():
            return False
        completed = threading.Event()
        answer = {"value": False}

        def ask():
            if self._cancel_requested():
                completed.set()
                return
            try:
                answer["value"] = messagebox.askokcancel(title, message)
            finally:
                completed.set()

        self.root.after(0, ask)
        completed.wait()
        return False if self._cancel_requested() else answer["value"]

    def log(self, message: str):
        logging.info(message)
        self.log_text.configure(state="normal")
        self.log_text.insert("end", f"  {message}\n")
        self.log_text.see("end")
        self.log_text.configure(state="disabled")

    def select_files(self, source_kind: str):
        if source_kind == "pdf":
            title = "添加 PDF 文件"
            filetypes = [("PDF 文件", "*.pdf")]
            extensions = {".pdf"}
            target_kind_after_selection = self._current_target_kind()
            if target_kind_after_selection == "pdf":
                target_kind_after_selection = PDF_TARGET_WORD
        elif source_kind == "word":
            title = "添加 Word 文件"
            filetypes = [("Word 文档", "*.docx")]
            extensions = {".docx"}
            target_kind_after_selection = "pdf"
        elif source_kind == "powerpoint":
            title = "添加 PowerPoint 文件"
            filetypes = [("PowerPoint 文件", "*.pptx *.ppt")]
            extensions = {".ppt", ".pptx"}
            target_kind_after_selection = "pdf"
        else:
            raise ValueError(f"Unsupported source kind: {source_kind}")

        selected = filedialog.askopenfilenames(title=title, filetypes=filetypes)
        if not selected:
            return 0

        selected_paths = [
            Path(path) for path in selected if Path(path).suffix.lower() in extensions
        ]
        return self._add_paths(
            source_kind,
            selected_paths,
            target_kind_after_selection,
        )

    def _add_paths(self, source_kind, selected_paths, target_kind):
        selected_paths = list(selected_paths)
        if not selected_paths:
            return 0

        if self.source_kind and self.source_kind != source_kind:
            confirmed = messagebox.askyesno(
                "切换转换方向",
                "当前队列包含另一种文件。是否清空现有队列并切换转换方向？",
            )
            if not confirmed:
                return 0
            self._clear_queue(log_change=False)

        if source_kind == "pdf" and not self._confirm_pdf_target_change(target_kind):
            return 0
        self._set_target_kind(target_kind)
        previous_count = len(self.input_paths)
        self.input_paths = deduplicate_paths(self.input_paths + selected_paths)
        self.source_kind = source_kind if self.input_paths else None
        self._refresh_queue()
        added_count = len(self.input_paths) - previous_count
        self.log(f"已添加 {added_count} 个文件；队列共 {len(self.input_paths)} 个")
        return added_count

    def _path_key(self, path: Path) -> str:
        return os.path.normcase(os.path.normpath(str(Path(path).absolute())))

    def _refresh_queue(self):
        self.file_tree.delete(*self.file_tree.get_children())
        self._tree_paths = {}
        for path in self.input_paths:
            item_id = self.file_tree.insert(
                "", "end",
                values=(path.name, str(path.parent), self.STATUS_TEXT["pending"], ""),
            )
            self._tree_paths[item_id] = path

        spec = self._current_spec()
        direction = spec.label.replace(" -> ", " → ") if spec else ""
        summary = f"{direction} · {len(self.input_paths)} 个文件" if direction else "尚未添加文件"
        self.queue_summary_var.set(summary)
        if self.input_paths:
            self.empty_drop_hint.place_forget()
        else:
            self.empty_drop_hint.place(relx=0.5, rely=0.56, anchor="center")
        self._update_method_panel()
        self._update_action_states()

    def remove_selected(self):
        selected_ids = self.file_tree.selection()
        selected_keys = {
            self._path_key(self._tree_paths[item_id])
            for item_id in selected_ids if item_id in self._tree_paths
        }
        if not selected_keys:
            return
        self.input_paths = [
            path for path in self.input_paths if self._path_key(path) not in selected_keys
        ]
        if not self.input_paths:
            self.source_kind = None
            if self._current_target_kind() == "pdf":
                self._set_target_kind(PDF_TARGET_WORD)
        self._refresh_queue()
        self.log(f"已移除 {len(selected_keys)} 个文件")

    def clear_queue(self):
        self._clear_queue(log_change=True)

    def _clear_queue(self, log_change):
        had_items = bool(self.input_paths)
        self.input_paths = []
        self.source_kind = None
        if self._current_target_kind() == "pdf":
            self._set_target_kind(PDF_TARGET_WORD)
        self._refresh_queue()
        if had_items and log_change:
            self.log("已清空转换队列")

    def choose_output_dir(self):
        initial_dir = self.output_dir_var.get() or None
        directory = filedialog.askdirectory(title="选择统一输出目录", initialdir=initial_dir)
        if directory:
            self.output_dir_var.set(directory)
            self.output_mode_var.set("custom")
            self._toggle_output_mode()

    def _toggle_output_mode(self):
        self._update_action_states()

    def _validate_output_dir(self):
        if self.output_mode_var.get() == "source":
            return None
        raw_path = self.output_dir_var.get().strip()
        if not raw_path:
            raise ValueError("请选择统一输出目录")
        output_dir = Path(raw_path)
        if not output_dir.is_dir():
            raise ValueError("输出目录不存在")
        try:
            import tempfile
            with tempfile.NamedTemporaryFile(dir=str(output_dir), prefix=".pdf_word_converter_"):
                pass
        except OSError as exc:
            raise ValueError(f"输出目录不可写: {exc}") from exc
        return output_dir

    def _set_preflight_status(self, current, total, source_name):
        self.progress_text_var.set(
            f"正在检查复杂版式 [{current}/{total}] {source_name}"
        )

    def _log_from_worker(self, message):
        self.root.after(0, self.log, message)

    def _best_editable_pdf_word_method(self, preferred_method):
        candidates = []
        is_macos = getattr(self, "is_macos", False)
        if preferred_method in (WORD_NATIVE, "libreoffice"):
            candidates.append(preferred_method)
        if not is_macos:
            candidates.append(WORD_NATIVE)
        candidates.append("libreoffice")

        for candidate in dict.fromkeys(candidates):
            if candidate == WORD_NATIVE and is_macos:
                continue
            if self._method_availability(candidate):
                return candidate
        return None

    @staticmethod
    def _unverifiable_table_form_reasons(report):
        reasons = []
        widget_count = int(getattr(report, "widget_count", 0) or 0)
        checkbox_symbol_count = int(
            getattr(report, "checkbox_symbol_count", 0) or 0
        )
        vector_mark_count = int(
            getattr(report, "vector_mark_count", 0) or 0
        )
        symbol_font_count = int(getattr(report, "symbol_font_run_count", 0) or 0)
        if widget_count:
            reasons.append(f"PDF 表单控件 {widget_count} 个")
        if checkbox_symbol_count:
            reasons.append(f"复选或勾选符号 {checkbox_symbol_count} 个")
        if vector_mark_count:
            reasons.append(f"疑似矢量复选或勾选标记 {vector_mark_count} 处")
        if symbol_font_count:
            reasons.append(f"符号字体文本 {symbol_font_count} 处")
        return reasons

    @staticmethod
    def _summarize_fidelity_reports(reports, reason_getter, limit=6):
        lines = []
        for report in reports[:limit]:
            reasons = reason_getter(report)
            detail = "；".join(reasons[:3]) if reasons else "复杂版式"
            lines.append(f"- {report.path.name}：{detail}")
        if len(reports) > limit:
            lines.append(f"- 另有 {len(reports) - limit} 个文件")
        return "\n".join(lines)

    def _choose_pdf_word_fidelity_mode(self, method, paths=None):
        """Build a per-file policy; editable table PDFs may never become images."""
        sources = list(paths if paths is not None else self.input_paths)
        policy = PdfWordBatchPolicy(default_method=method)

        def abandon_policy():
            policy.cleanup_snapshots()
            return None

        editable_reports = []
        protected_text_reports = []
        unconfirmed_table_reports = []
        unverifiable_table_reports = []
        image_candidate_reports = []
        checked_files = 0
        failed_files = 0

        def add_completion_warning(source_path, message):
            prefix = f"{source_path.name}："
            message = str(message)
            if message.startswith(prefix):
                message = message[len(prefix):]
            existing = policy.completion_warnings.get(source_path, "")
            if existing and message not in existing:
                policy.completion_warnings[source_path] = f"{existing}；{message}"
            elif not existing:
                policy.completion_warnings[source_path] = f"{prefix}{message}"

        for index, source in enumerate(sources, start=1):
            if self._cancel_requested():
                return abandon_policy()
            source_path = Path(source)
            self.root.after(
                0,
                self._set_preflight_status,
                index,
                len(sources),
                source_path.name,
            )
            snapshot_directory = None
            try:
                identity_getter = getattr(
                    self,
                    "_source_file_identity",
                    _source_file_identity,
                )
                snapshot_copy = getattr(
                    self,
                    "_copy_verified_source_snapshot",
                    _copy_verified_source_snapshot,
                )
                identity_before = identity_getter(source_path)
                snapshot_directory = tempfile.TemporaryDirectory(
                    prefix="pdf-word-converter-preflight-"
                )
                snapshot_path = (
                    Path(snapshot_directory.name) / source_path.name
                )
                snapshot_copy(
                    source_path,
                    snapshot_path,
                    identity_before,
                )
                report = analyze_pdf_fidelity_risk(
                    snapshot_path,
                    max_pages=None,
                    cancel_requested=self._cancel_requested,
                )
                if isinstance(report, PdfFidelityRisk):
                    report = replace(report, path=source_path)
                else:
                    # Compatible analyzer doubles may expose a mutable path
                    # attribute without being dataclasses.
                    report.path = source_path
                identity_after = identity_getter(source_path)
                if identity_before != identity_after:
                    raise RuntimeError(
                        "文件在版式预检期间发生变化，请重新加入后再转换"
                    )
                policy.source_identities[source_path] = identity_after
                policy.source_snapshots[source_path] = snapshot_path
                policy.snapshot_directories.append(snapshot_directory)
                snapshot_directory = None
            except PdfFidelityAnalysisCancelled:
                if snapshot_directory is not None:
                    snapshot_directory.cleanup()
                self._log_from_worker("版式预检已取消")
                return abandon_policy()
            except Exception as exc:
                if snapshot_directory is not None:
                    snapshot_directory.cleanup()
                failed_files += 1
                policy.unverified_paths.add(source_path)
                policy.image_protected_paths.add(source_path)
                self._log_from_worker(
                    f"版式预检跳过 {source_path.name}: {exc}"
                )
                policy.blocked_paths[source_path] = (
                    "无法完成该 PDF 的可编辑内容预检；为避免把可能存在的可编辑"
                    "表格或文字转成照片，本项已停止"
                )
                continue

            checked_files += 1
            table_uncertain = bool(
                getattr(report, "table_analysis_uncertain", False)
                or (
                    int(getattr(report, "table_count", 0) or 0)
                    and int(
                        getattr(
                            report,
                            "table_text_extraction_failure_count",
                            0,
                        )
                        or 0
                    )
                )
            )
            has_selectable_text = bool(
                getattr(report, "has_selectable_text", False)
                or int(getattr(report, "selectable_text_character_count", 0) or 0)
            )
            editable_table = bool(
                getattr(report, "editable_table_candidate", False)
                or table_uncertain
            )
            unconfirmed_table_regions = int(
                getattr(report, "unconfirmed_table_region_count", 0) or 0
            )
            unconfirmed_table_layout = bool(
                unconfirmed_table_regions
                or (
                    getattr(report, "table_layout_suspected", False)
                    and not int(getattr(report, "table_count", 0) or 0)
                )
            )
            if editable_table or has_selectable_text:
                policy.image_protected_paths.add(source_path)

            mapping_warning_count = int(
                getattr(report, "unverifiable_text_character_count", 0) or 0
            )
            if mapping_warning_count:
                add_completion_warning(
                    source_path,
                    f"{source_path.name}：检测到 {mapping_warning_count} 个无法可靠"
                    "映射的文字字符，可能来自损坏字体编码或私用区字符。转换将继续，"
                    "完成后请对照原 PDF 人工核对这些字符",
                )

            if bool(getattr(report, "table_analysis_limited", False)):
                add_completion_warning(
                    source_path,
                    f"{source_path.name}：文档页数较多，已快速检查全部页面文字，"
                    "但未逐页重建并严格核对所有表格；转换完成后请重点检查表格行列、"
                    "边框和勾选状态",
                )

            if table_uncertain:
                policy.advisory_validation_paths.add(source_path)
                uncertain_count = int(
                    getattr(report, "table_text_extraction_failure_count", 0) or 0
                )
                uncertain_label = (
                    f"{uncertain_count} 个" if uncertain_count else "部分"
                )
                add_completion_warning(
                    source_path,
                    f"{source_path.name}：有 {uncertain_label}源表格无法"
                    "建立完整的逐格校验基线；转换仍会继续，完成后请人工核对这些表格",
                )

            form_reasons = self._unverifiable_table_form_reasons(report)
            if form_reasons:
                policy.advisory_validation_paths.add(source_path)
                add_completion_warning(
                    source_path,
                    f"{source_path.name}：检测到无法自动逐项验证的"
                    + "、".join(form_reasons)
                    + "；转换仍会继续，完成后请重点核对对应符号和勾选状态",
                )

            if has_selectable_text:
                policy.selectable_text_reports[source_path] = report

            if unconfirmed_table_layout:
                unconfirmed_table_reports.append(report)
                policy.advisory_validation_paths.add(source_path)
                add_completion_warning(
                    source_path,
                    f"{source_path.name}：发现 {unconfirmed_table_regions or '若干'} 处"
                    "疑似表格区域，但无法建立可靠的逐格校验基线；已优先继续转换，"
                    "完成后请对照原 PDF 仔细检查",
                )

            if editable_table:
                if form_reasons:
                    unverifiable_table_reports.append(report)
                editable_reports.append(report)
                policy.editable_table_reports[source_path] = report
                editable_method = self._best_editable_pdf_word_method(method)
                if editable_method is None:
                    policy.blocked_paths[source_path] = (
                        "检测到可编辑表格或表格结构无法完全验证，但当前没有可用的 Word 或 "
                        "LibreOffice 可编辑转换引擎；为避免生成照片，本项未转换"
                    )
                else:
                    policy.method_overrides[source_path] = editable_method
            elif has_selectable_text:
                protected_text_reports.append(report)
                editable_method = self._best_editable_pdf_word_method(method)
                if editable_method is None:
                    policy.blocked_paths[source_path] = (
                        "检测到可选择文字，但当前没有可用的 Word 或 LibreOffice "
                        "可编辑转换引擎；为避免遗漏表格并生成照片，本项未转换"
                    )
                else:
                    policy.method_overrides[source_path] = editable_method
            elif report.is_complex:
                image_candidate_reports.append(report)

        if self._cancel_requested():
            return abandon_policy()

        if editable_reports:
            table_text_count = sum(
                int(getattr(report, "table_text_character_count", 0) or 0)
                for report in editable_reports
            )
            detail = f"；另有 {failed_files} 个无法检查" if failed_files else ""
            self._log_from_worker(
                f"检测到 {len(editable_reports)} 个含可选择文字的表格 PDF"
                f"（表格文字共 {table_text_count} 字符），已锁定可编辑转换并启用"
                f"输出表格结构校验{detail}"
            )

        if protected_text_reports:
            self._log_from_worker(
                f"检测到 {len(protected_text_reports)} 个含可选择文字但未确认表格的 PDF，"
                "已禁止整页图片转换和图片回退"
            )

        if unconfirmed_table_reports:
            self._log_from_worker(
                f"检测到 {len(unconfirmed_table_reports)} 个版面疑似表格但无法提取"
                "可验证结构的 PDF，已继续使用可编辑引擎并将在完成后提醒人工核对"
            )

        if unverifiable_table_reports:
            details = self._summarize_fidelity_reports(
                unverifiable_table_reports,
                self._unverifiable_table_form_reasons,
            )
            self._log_from_worker(
                "检测到无法可靠验证勾选状态或符号字体的可编辑表格，"
                f"已继续转换并将在完成后提醒人工核对：{details}"
            )

        forced_non_image_reports = [
            report
            for report in editable_reports + protected_text_reports
            if policy.method_overrides.get(Path(report.path)) not in (None, method)
        ]
        if forced_non_image_reports:
            if self._cancel_requested():
                return abandon_policy()
            details = self._summarize_fidelity_reports(
                forced_non_image_reports,
                lambda report: [
                    f"可选择文字 {getattr(report, 'selectable_text_character_count', 0)} 字符",
                    f"表格结构 {getattr(report, 'table_count', 0)} 个",
                    f"表格文字 {getattr(report, 'table_text_character_count', 0)} 字符",
                    f"将使用 {policy.method_overrides[Path(report.path)]}",
                ],
            )
            use_editable = self._ask_ok_cancel_from_worker(
                "可选择文字不能使用图片模式",
                "你选择了整页图片模式，但以下 PDF 含可选择文字，可能包含可编辑表格：\n"
                f"{details}\n\n"
                "为保证文字和可能存在的表格可编辑，程序将只对这些文件改用可编辑引擎；"
                "其他文件仍按图片模式转换。\n\n"
                "选择“确定”继续；选择“取消”停止本批。",
            )
            if not use_editable:
                self._log_from_worker("用户取消了图片模式中的可编辑内容转换")
                return abandon_policy()

        blocked_count = len(policy.blocked_paths)
        if blocked_count:
            blocked_names = "\n".join(
                f"- {path.name}" for path in list(policy.blocked_paths)[:6]
            )
            if blocked_count > 6:
                blocked_names += f"\n- 另有 {blocked_count - 6} 个文件"
            if self._cancel_requested():
                return abandon_policy()
            continue_batch = self._ask_ok_cancel_from_worker(
                "部分文件无法保证可编辑内容",
                f"以下 {blocked_count} 个文件没有可用的安全转换路径：\n"
                f"{blocked_names}\n\n"
                "这些文件会标记为失败且不会生成照片；批次中的其他文件仍可继续。\n\n"
                "选择“确定”继续；选择“取消”停止本批。",
            )
            if not continue_batch:
                self._log_from_worker("用户取消了缺少可编辑引擎的转换")
                return abandon_policy()

        if method != "images" and image_candidate_reports:
            if self._cancel_requested():
                return abandon_policy()
            details = self._summarize_fidelity_reports(
                image_candidate_reports,
                lambda report: list(getattr(report, "reasons", ())),
            )
            if not self._method_availability("images"):
                continue_editable = self._ask_ok_cancel_from_worker(
                    "检测到扫描或纯图形复杂版式",
                    f"以下文件没有可选择文字，但存在复杂图形：\n{details}\n\n"
                    "当前内置图片组件不可用，只能继续可编辑模式，版式可能变化。\n\n"
                    "选择“确定”继续；选择“取消”停止本批。",
                )
                if not continue_editable:
                    self._log_from_worker("用户取消了复杂版式转换")
                    return abandon_policy()
            else:
                choice = self._ask_yes_no_cancel_from_worker(
                    "检测到扫描或纯图形复杂版式",
                    f"以下文件没有可选择文字，但存在复杂图形：\n{details}\n\n"
                    "选择“是”：仅这些文件改用整页高清图片；\n"
                    "选择“否”：这些文件也继续可编辑模式；\n"
                    "选择“取消”：停止本批。\n\n"
                    "同批中的可编辑表格文件不会受到影响，也不会被改成图片。",
                )
                if choice is None:
                    self._log_from_worker("用户取消了复杂版式转换")
                    return abandon_policy()
                if choice:
                    for report in image_candidate_reports:
                        policy.method_overrides[Path(report.path)] = "images"
                    self._log_from_worker(
                        f"仅将 {len(image_candidate_reports)} 个非表格复杂 PDF "
                        "切换为整页图片模式"
                    )
                else:
                    self._log_from_worker(
                        f"{len(image_candidate_reports)} 个非表格复杂 PDF "
                        "继续使用可编辑模式"
                    )
        elif (
            not editable_reports
            and not protected_text_reports
            and not unconfirmed_table_reports
            and not image_candidate_reports
        ):
            if checked_files:
                detail = f"，另有 {failed_files} 个无法检查" if failed_files else ""
                self._log_from_worker(
                    f"版式预检完成：成功检查 {checked_files} 个 PDF，"
                    f"未发现高风险结构{detail}"
                )
            elif failed_files:
                self._log_from_worker(
                    f"版式预检未完成：{failed_files} 个 PDF 均无法检查，"
                    "已按安全规则标记失败"
                )

        return policy

    def start_conversion(self):
        if not self.input_paths:
            messagebox.showwarning("提示", "请先添加文件")
            return

        spec = self._current_spec()
        if spec is None:
            messagebox.showerror("队列错误", "无法确定当前转换方向")
            return
        if not self._conversion_detection_ready():
            messagebox.showinfo("提示", "当前转换方向所需引擎仍在检测，请稍候")
            return

        if spec.key == "pdf_to_word":
            method = self.method_var.get()
            if not self._method_availability(method):
                if method == WORD_NATIVE:
                    self._prompt_word_install()
                    return
                if method == "libreoffice":
                    self._prompt_libreoffice_install()
                    return
                messagebox.showerror("转换方式不可用", "请选择一个当前可用的 PDF 转换方式")
                return
            if self.is_macos and method == WORD_NATIVE:
                messagebox.showerror(
                    "此平台不支持",
                    "macOS 源码预览尚未开放 Word 原生 PDF → DOCX。含可选择文字或表格的 "
                    "PDF 请使用 LibreOffice；只有完全没有可选择文字的扫描件或纯图形 "
                    "PDF 才能使用图片模式。",
                )
                return
        elif spec.key == "pdf_to_powerpoint":
            method = "ppt_images"
            missing = []
            if not self._method_availability("images"):
                missing.append("PyMuPDF")
            if not self._method_availability("pptx"):
                missing.append("python-pptx")
            if missing:
                messagebox.showerror(
                    "内置组件不可用",
                    "PDF → PowerPoint 所需内置组件不可用：" + "、".join(missing)
                    + "。请重新下载安装完整程序。",
                )
                return
        elif spec.key == "word_to_pdf":
            method = None
            if not (self._method_availability(WORD_NATIVE) or self._method_availability("libreoffice")):
                if self.is_macos:
                    self._prompt_macos_conversion_engines(
                        "Microsoft Word",
                        WORD_NATIVE,
                        getattr(self, "_word_registered", False),
                    )
                else:
                    self._prompt_libreoffice_install()
                return
        elif spec.key == "powerpoint_to_pdf":
            method = None
            if not (
                self._method_availability(POWERPOINT_NATIVE)
                or self._method_availability("libreoffice")
            ):
                if self.is_macos:
                    self._prompt_macos_conversion_engines(
                        "Microsoft PowerPoint",
                        POWERPOINT_NATIVE,
                        getattr(self, "_powerpoint_registered", False),
                    )
                else:
                    self._prompt_libreoffice_install()
                return

        try:
            output_dir = self._validate_output_dir()
        except ValueError as exc:
            messagebox.showerror("输出目录错误", str(exc))
            return

        self.output_paths = []
        self._stop_cancel_wait_timer()
        self._cancel_event.clear()
        self._reset_queue_status()
        self._set_busy(True)
        self.progress["value"] = 0
        self.progress_text_var.set("准备转换...")
        direction = spec.label.replace(" -> ", " → ")
        self.log(f"开始 {direction} 批量转换，共 {len(self.input_paths)} 个文件")

        paths = list(self.input_paths)
        threading.Thread(
            target=self._prepare_and_run_batch,
            args=(paths, spec, output_dir, method),
            daemon=True,
        ).start()

    def _prepare_and_run_batch(self, paths, spec, output_dir, method):
        policy = None
        try:
            if spec.key == "pdf_to_word":
                policy = self._choose_pdf_word_fidelity_mode(method, paths)
                if policy is None:
                    self._cancel_event.set()
        except Exception as exc:
            self.root.after(0, self._on_batch_error, str(exc))
            return
        self._run_batch(paths, spec, output_dir, method, policy)

    def _reset_queue_status(self):
        for item_id in self.file_tree.get_children():
            values = list(self.file_tree.item(item_id, "values"))
            values[2] = self.STATUS_TEXT["pending"]
            values[3] = ""
            self.file_tree.item(item_id, values=values, tags=())

    @staticmethod
    def _validate_selectable_text_output(output, report, progress):
        summary = inspect_editable_docx_tables(output)
        source_text = getattr(report, "selectable_text", None)
        output_text = getattr(summary, "visible_text", None)
        source_count = int(
            getattr(report, "selectable_text_character_count", 0) or 0
        )
        output_count = int(
            getattr(summary, "visible_text_character_count", 0) or 0
        )
        mapping_warning_count = int(
            getattr(report, "unverifiable_text_character_count", 0) or 0
        )
        if mapping_warning_count:
            progress(
                f"检测到 {mapping_warning_count} 个字体映射可疑字符，"
                "已保留结果并将在完成时提醒人工核对",
                99,
            )
            return
        if not _document_text_evidence_matches(
            source_text,
            output_text,
            source_character_count=source_count,
            output_character_count=output_count,
        ):
            raise RuntimeError(
                "Word 输出未通过全文文字校验：可见文字 "
                f"{output_count}/{source_count} 个非空字符。"
                "为避免交付漏字、错字或额外文字的文件，结果已拒绝"
            )
        progress(f"已确认全文可见文字：{output_count} 个非空字符", 99)

    @staticmethod
    def _validate_editable_table_output(output, report, progress):
        summary = inspect_editable_docx_tables(output)
        source_text_count = int(
            getattr(report, "table_text_character_count", 0) or 0
        )
        source_table_count = max(
            1,
            int(getattr(report, "table_count", 0) or 0),
        )
        minimum_text_count = source_text_count
        source_table_texts = tuple(
            getattr(report, "table_texts", ()) or ()
        )
        source_table_cell_matrices = tuple(
            getattr(report, "table_cell_matrices", ()) or ()
        )
        source_table_shapes = tuple(
            getattr(report, "table_shapes", ()) or ()
        )
        source_table_full_width_span_rows = tuple(
            getattr(report, "table_full_width_span_rows", ()) or ()
        )
        if not source_table_full_width_span_rows:
            source_table_full_width_span_rows = tuple(
                (False,) * int(shape[0])
                for shape in source_table_shapes
                if _coerce_table_shape(shape) is not None
            )
        source_table_cell_counts = tuple(
            getattr(report, "table_cell_counts", ()) or ()
        )
        source_table_cell_spans = tuple(
            getattr(report, "table_cell_spans", ()) or ()
        )
        source_table_cell_border_edges = tuple(
            getattr(report, "table_cell_border_edges", ()) or ()
        )
        output_table_texts = tuple(
            getattr(summary, "table_texts", ()) or ()
        )
        output_table_cell_matrices = tuple(
            getattr(summary, "table_cell_matrices", ()) or ()
        )
        output_table_shapes = tuple(
            getattr(summary, "table_shapes", ()) or ()
        )
        output_table_cell_counts = tuple(
            getattr(summary, "table_cell_counts", ()) or ()
        )
        output_table_cell_spans = tuple(
            getattr(summary, "table_cell_spans", ()) or ()
        )
        output_table_cell_border_edges = tuple(
            getattr(summary, "table_cell_border_edges", ()) or ()
        )
        source_document_text = getattr(report, "selectable_text", None)
        output_document_text = getattr(summary, "visible_text", None)
        source_document_text_count = int(
            getattr(report, "selectable_text_character_count", 0) or 0
        )
        output_document_text_count = int(
            getattr(summary, "visible_text_character_count", 0) or 0
        )
        mapping_warning_count = int(
            getattr(report, "unverifiable_text_character_count", 0) or 0
        )
        document_text_evidence_valid = (
            True
            if source_document_text is None or mapping_warning_count
            else _document_text_evidence_matches(
                source_document_text,
                output_document_text,
                source_character_count=source_document_text_count,
                output_character_count=output_document_text_count,
            )
        )
        source_table_border_maps = tuple(
            _coerce_table_cell_border_map(edges, spans, shape)
            for edges, spans, shape in zip(
                source_table_cell_border_edges,
                source_table_cell_spans,
                source_table_shapes,
            )
        )
        output_table_border_maps = tuple(
            _coerce_table_cell_border_map(edges, spans, shape)
            for edges, spans, shape in zip(
                output_table_cell_border_edges,
                output_table_cell_spans,
                output_table_shapes,
            )
        )
        border_evidence_valid = (
            len(source_table_cell_border_edges) == source_table_count
            and len(output_table_cell_border_edges) == summary.table_count
            and len(source_table_border_maps) == source_table_count
            and len(output_table_border_maps) == summary.table_count
            and all(border_map is not None for border_map in source_table_border_maps)
            and all(border_map is not None for border_map in output_table_border_maps)
            and source_table_border_maps == output_table_border_maps
        )
        source_structure_complete = (
            len(source_table_texts) == source_table_count
            and len(source_table_cell_matrices) == source_table_count
            and len(source_table_shapes) == source_table_count
            and len(source_table_full_width_span_rows) == source_table_count
            and len(source_table_cell_counts) == source_table_count
            and len(source_table_cell_spans) == source_table_count
            and len(source_table_cell_border_edges) == source_table_count
            and len(source_table_border_maps) == source_table_count
            and all(border_map is not None for border_map in source_table_border_maps)
            and all(
                _coerce_table_cell_spans(spans, shape) is not None
                and len(spans) == int(cell_count)
                for spans, shape, cell_count in zip(
                    source_table_cell_spans,
                    source_table_shapes,
                    source_table_cell_counts,
                )
            )
            and all(
                len(table_rows) == int(shape[0])
                for table_rows, shape in zip(
                    source_table_full_width_span_rows,
                    source_table_shapes,
                )
            )
            and all(
                _table_cell_matrix_matches_text(matrix, shape, text)
                for matrix, shape, text in zip(
                    source_table_cell_matrices,
                    source_table_shapes,
                    source_table_texts,
                )
            )
        )
        output_structure_complete = (
            len(output_table_texts) == summary.table_count
            and len(output_table_cell_matrices) == summary.table_count
            and len(output_table_shapes) == summary.table_count
            and len(output_table_cell_counts) == summary.table_count
            and len(output_table_cell_spans) == summary.table_count
            and len(output_table_cell_border_edges) == summary.table_count
            and len(output_table_border_maps) == summary.table_count
            and all(border_map is not None for border_map in output_table_border_maps)
            and all(
                _coerce_table_cell_spans(spans, shape) is not None
                and len(spans) == int(cell_count)
                for spans, shape, cell_count in zip(
                    output_table_cell_spans,
                    output_table_shapes,
                    output_table_cell_counts,
                )
            )
            and all(
                _table_cell_matrix_matches_text(matrix, shape, text)
                for matrix, shape, text in zip(
                    output_table_cell_matrices,
                    output_table_shapes,
                    output_table_texts,
                )
            )
        )
        table_content_valid, matched_table_count, compared_table_count = (
            _table_structures_preserved(
                source_table_texts,
                source_table_shapes,
                source_table_cell_counts,
                output_table_texts,
                output_table_shapes,
                output_table_cell_counts,
                source_cell_matrices=source_table_cell_matrices,
                output_cell_matrices=output_table_cell_matrices,
                source_full_width_span_rows=(
                    source_table_full_width_span_rows
                ),
                source_cell_spans=source_table_cell_spans,
                output_cell_spans=output_table_cell_spans,
            )
        )
        source_shape_text = "、".join(
            f"{shape[0]}×{shape[1]}/{cell_count} 格"
            for shape, cell_count in zip(
                source_table_shapes,
                source_table_cell_counts,
            )
        ) or "缺失"
        output_shape_text = "、".join(
            f"{shape[0]}×{shape[1]}/{cell_count} 格"
            for shape, cell_count in zip(
                output_table_shapes,
                output_table_cell_counts,
            )
        ) or "缺失"
        has_editable_structure = bool(
            getattr(
                summary,
                "has_editable_table_structure",
                summary.has_editable_table,
            )
        )
        has_required_editable_table = (
            bool(summary.has_editable_table)
            if minimum_text_count
            else has_editable_structure
        )
        large_page_drawing_count = int(
            getattr(summary, "large_page_drawing_count", 0) or 0
        )
        if (
            not has_required_editable_table
            or summary.table_count != source_table_count
            or summary.table_text_character_count != minimum_text_count
            or large_page_drawing_count
            or not source_structure_complete
            or not output_structure_complete
            or not table_content_valid
            or not border_evidence_valid
            or not document_text_evidence_valid
        ):
            protection_details = []
            if not document_text_evidence_valid:
                protection_details.append(
                    "整份文档的可见文字与源 PDF 不一致"
                )
            if not border_evidence_valid:
                protection_details.append("单元格边框与源 PDF 不一致")
            if summary.document_protected:
                protection_details.append("文档被限制编辑")
            locked_tables = int(
                getattr(summary, "locked_content_control_table_count", 0) or 0
            )
            if locked_tables:
                protection_details.append(f"{locked_tables} 个表格被内容控件锁定")
            if large_page_drawing_count:
                protection_details.append(
                    f"检测到 {large_page_drawing_count} 张大面积页面图片"
                )
            invalid_layout_tables = int(
                getattr(summary, "invalid_layout_table_count", 0) or 0
            )
            if invalid_layout_tables:
                protection_details.append(
                    f"{invalid_layout_tables} 个表格无法证明在页面中清晰可见"
                )
            protection_detail = (
                "，且" + "、".join(protection_details)
                if protection_details
                else ""
            )
            raise RuntimeError(
                "源 PDF 含可编辑表格结构，但输出未通过可编辑表格校验："
                f"Word 表格 {summary.table_count}/{source_table_count} 个、"
                f"行 {summary.row_count} 个、"
                f"单元格 {summary.cell_count} 个、表格文字 "
                f"{summary.table_text_character_count}/{minimum_text_count} 字符、"
                f"逐表、逐行和逐格内容与结构 {matched_table_count}/{compared_table_count} 个"
                f"（源 {source_shape_text}；输出 {output_shape_text}）、"
                f"全文可见文字 {output_document_text_count}/"
                f"{source_document_text_count} 个非空字符"
                f"{protection_detail}。为避免交付照片、伪表格或漏字结果，结果已拒绝"
            )
        progress(
            f"已确认可编辑 Word 表格：{summary.table_count} 个表格，"
            f"{summary.table_text_character_count} 个表格文字字符，"
            f"全文 {output_document_text_count} 个非空字符",
            99,
        )

    def _run_batch(self, paths, spec, output_dir, method, policy=None):
        completion_warnings = {}
        if spec.key != "pdf_to_word":
            converter = self._create_converter(spec, method)
        else:
            if policy is None:
                if self._cancel_requested():
                    policy = PdfWordBatchPolicy(default_method=method)
                else:
                    self.root.after(
                        0,
                        self._on_batch_error,
                        "PDF → Word 缺少版式预检结果，已阻止转换",
                    )
                    return
            method_overrides = {
                self._path_key(path): item_method
                for path, item_method in policy.method_overrides.items()
            }
            editable_reports = {
                self._path_key(path): report
                for path, report in policy.editable_table_reports.items()
            }
            selectable_text_reports = {
                self._path_key(path): report
                for path, report in policy.selectable_text_reports.items()
            }
            blocked_paths = {
                self._path_key(path): message
                for path, message in policy.blocked_paths.items()
            }
            unverified_paths = {
                self._path_key(path) for path in policy.unverified_paths
            }
            image_protected_paths = {
                self._path_key(path) for path in policy.image_protected_paths
            }
            source_identities = {
                self._path_key(path): identity
                for path, identity in policy.source_identities.items()
            }
            source_snapshots = {
                self._path_key(path): Path(snapshot)
                for path, snapshot in policy.source_snapshots.items()
            }
            completion_warnings = {
                self._path_key(path): message
                for path, message in policy.completion_warnings.items()
            }
            advisory_validation_paths = {
                self._path_key(path) for path in policy.advisory_validation_paths
            }
            converter_cache = {}

            def converter(source, output, progress):
                source_path = Path(source)
                source_key = self._path_key(source_path)
                blocked_message = blocked_paths.get(source_key)
                if blocked_message:
                    raise RuntimeError(blocked_message)

                expected_identity = source_identities.get(source_key)
                if expected_identity is None:
                    raise RuntimeError(
                        "PDF 缺少与版式预检绑定的源文件身份，本项已停止，"
                        "请重新加入文件后再转换"
                    )
                identity_getter = getattr(
                    self,
                    "_source_file_identity",
                    _source_file_identity,
                )
                if identity_getter(source_path) != expected_identity:
                    raise RuntimeError(
                        "源 PDF 在版式预检后发生变化；为避免沿用旧判断，本项已停止，"
                        "请重新加入文件后再转换"
                    )

                item_method = method_overrides.get(
                    source_key,
                    policy.default_method,
                )
                report = editable_reports.get(source_key)
                selectable_text_report = selectable_text_reports.get(source_key)
                prevent_image_fallback = (
                    report is not None
                    or source_key in unverified_paths
                    or source_key in image_protected_paths
                )
                if item_method == "images" and prevent_image_fallback:
                    raise RuntimeError(
                        "该 PDF 含可选择文字或可能含可编辑表格，已阻止整页图片转换"
                    )

                cache_key = (item_method, prevent_image_fallback)
                item_converter = converter_cache.get(cache_key)
                if item_converter is None:
                    item_converter = self._create_converter(
                        spec,
                        item_method,
                        allow_image_fallback=not prevent_image_fallback,
                    )
                    converter_cache[cache_key] = item_converter

                snapshot_path = source_snapshots.get(source_key)
                if snapshot_path is None or not snapshot_path.is_file():
                    raise RuntimeError(
                        "PDF 预检快照不存在，本项已停止，请重新加入文件后再转换"
                    )
                if self._path_key(snapshot_path) == source_key:
                    raise RuntimeError("PDF 预检快照不能与原始文件使用同一路径")
                if _source_file_identity(snapshot_path)[4] != expected_identity[4]:
                    raise RuntimeError("PDF 预检快照与源文件身份不一致")
                engine_progress = progress
                if report is not None or selectable_text_report is not None:
                    def engine_progress(message, pct):
                        bounded_pct = max(0, min(100, int(pct)))
                        progress(message, min(96, int(bounded_pct * 0.96)))
                item_converter(str(snapshot_path), output, engine_progress)
                if identity_getter(source_path) != expected_identity:
                    raise RuntimeError(
                        "源 PDF 在转换期间发生变化；本项已停止，"
                        "请重新加入文件后再转换"
                    )
                if _source_file_identity(snapshot_path)[4] != expected_identity[4]:
                    raise RuntimeError("转换引擎意外修改了受保护的源 PDF 快照")
                if report is not None:
                    progress("正在校正可编辑 Word 表格结构...", 97)
                    mapping_warning_count = int(
                        getattr(
                            report,
                            "unverifiable_text_character_count",
                            0,
                        )
                        or 0
                    )
                    advisory_validation = bool(
                        mapping_warning_count
                        or source_key in advisory_validation_paths
                    )
                    if advisory_validation:
                        # An advisory result must still be a readable DOCX.
                        inspect_editable_docx_tables(output)
                    try:
                        repaired = repair_docx_table_topology(
                            output,
                            table_shapes=report.table_shapes,
                            table_cell_matrices=report.table_cell_matrices,
                            table_cell_spans=report.table_cell_spans,
                            table_cell_border_edges=report.table_cell_border_edges,
                            table_column_widths=report.table_column_widths,
                        )
                        if repaired:
                            progress("已安全恢复 Word 表格的行列与合并关系", 98)
                        self._validate_editable_table_output(
                            output,
                            report,
                            progress,
                        )
                    except Exception as exc:
                        if not advisory_validation:
                            if isinstance(exc, DocxTableRepairError):
                                raise RuntimeError(
                                    "源 PDF 含可编辑表格，但 Word 输出无法在不猜测的情况下"
                                    f"恢复原始单元格结构：{exc}"
                                ) from exc
                            raise
                        original_warning = completion_warnings.get(source_key, "")
                        validation_warning = (
                            "可编辑表格的严格文字或结构校验未完全通过："
                            f"{exc}。已按你的设置保留可打开的转换结果，请重点人工核对"
                        )
                        completion_warnings[source_key] = (
                            f"{original_warning}；{validation_warning}"
                            if original_warning
                            else f"{source_path.name}：{validation_warning}"
                        )
                        progress(
                            "部分结构无法严格验证；"
                            "已保留结果并将在完成时提醒",
                            99,
                        )
                elif selectable_text_report is not None:
                    progress("正在校验 Word 全文可见文字...", 97)
                    try:
                        self._validate_selectable_text_output(
                            output,
                            selectable_text_report,
                            progress,
                        )
                    except Exception as exc:
                        if source_key not in advisory_validation_paths:
                            raise
                        original_warning = completion_warnings.get(source_key, "")
                        validation_warning = (
                            "全文文字严格校验未完全通过："
                            f"{exc}。已保留可打开的转换结果，请对照原 PDF 仔细检查"
                        )
                        completion_warnings[source_key] = (
                            f"{original_warning}；{validation_warning}"
                            if original_warning
                            else f"{source_path.name}：{validation_warning}"
                        )
                        progress(
                            "全文文字无法完全一致；已保留结果并将在完成时提醒",
                            99,
                        )
        try:
            results = run_conversion_batch(
                paths,
                spec.target_suffix,
                output_dir,
                converter,
                self._cancel_event,
                self._progress_update,
                self._status_update,
            )
            successful_keys = {
                self._path_key(result.source)
                for result in results
                if result.status == "success"
            }
            warning_messages = [
                message
                for path_key, message in completion_warnings.items()
                if path_key in successful_keys
            ]
            self.root.after(
                0,
                self._on_batch_complete,
                results,
                warning_messages,
            )
        except Exception as exc:
            self.root.after(0, self._on_batch_error, str(exc))
        finally:
            if policy is not None:
                policy.cleanup_snapshots()

    def _record_native_status(self, key, status):
        self._set_engine_status(key, status)
        self.root.after(0, self._update_engine_indicators)
        self.root.after(0, self._refresh_engine_status_text)
        self.root.after(0, self._update_action_states)

    def _convert_with_macos_backend(self, key, backend, source, output, progress):
        try:
            backend.convert(source, output, progress)
        except MacOSOfficeError as exc:
            self._record_native_status(key, exc.status)
            raise
        else:
            self._record_native_status(
                key,
                EngineStatus(EngineState.AVAILABLE, "转换已验证"),
            )

    def _convert_with_windows_fallback(
        self,
        native_name,
        native_converter,
        libreoffice_converter,
        source,
        output,
        progress,
    ):
        try:
            native_converter(source, output, progress)
        except Exception as native_error:
            if not _engine_result_selectable(libreoffice_available()):
                raise
            output_path = Path(output)
            if output_path.exists():
                output_path.unlink()
            progress(f"{native_name} 原生转换失败，改用 LibreOffice 重试...", 5)
            try:
                libreoffice_converter(source, output, progress)
            except Exception as fallback_error:
                raise RuntimeError(
                    f"{native_name} 原生转换失败：{native_error}\n"
                    f"LibreOffice 回退也失败：{fallback_error}"
                ) from fallback_error

    def _create_converter(self, spec, method, allow_image_fallback=True):
        if spec.key == "pdf_to_word":
            def convert_pdf(source, output, progress):
                if method == WORD_NATIVE and not self.is_macos:
                    progress("使用 Microsoft Word 可编辑模式...", 1)
                    pdf_to_word_via_word(source, output, progress)
                elif method == "images":
                    progress("使用外观优先的整页高清图片模式...", 1)
                    pdf_to_word_via_images(source, output, progress)
                elif method == "libreoffice":
                    try:
                        pdf_to_word_via_libreoffice(
                            source, output,
                            lambda message, pct: progress(message, min(88, int(pct * 0.88))),
                        )
                    except Exception as exc:
                        if not allow_image_fallback:
                            raise RuntimeError(
                                "可编辑表格保护已启用，LibreOffice 转换失败后未回退到"
                                f"整页图片：{exc}"
                            ) from exc
                        can_retry = pymupdf_available() and is_libreoffice_export_filter_error(exc)
                        if not can_retry or not self._ask_yes_no_from_worker(
                            "LibreOffice 转换失败",
                            "LibreOffice 对 PDF → Word 的兼容性有限，当前文件无法导出。\n\n"
                            "是否改用内置图片模式重试？图片模式能保留页面观感，但文字不可编辑。",
                        ):
                            raise
                        output_path = Path(output)
                        if output_path.exists():
                            output_path.unlink()
                        progress("改用内置图片模式重试...", 5)
                        pdf_to_word_via_images(source, output, progress)
                        return
                else:
                    raise RuntimeError(f"未知转换方式: {method}")

            return convert_pdf

        if spec.key == "pdf_to_powerpoint":
            return pdf_to_pptx_via_images

        if spec.key == "word_to_pdf":
            def convert_docx(source, output, progress):
                if self.is_macos and self._method_availability(WORD_NATIVE):
                    try:
                        self._convert_with_macos_backend(
                            WORD_NATIVE,
                            self._mac_word_backend,
                            source,
                            output,
                            progress,
                        )
                    except MacOSOfficeError as native_error:
                        if not _engine_result_selectable(libreoffice_available()):
                            raise
                        output_path = Path(output)
                        if output_path.exists():
                            output_path.unlink()
                        progress(
                            "Word 原生转换失败，改用 LibreOffice 重试（详见自动化权限状态）...",
                            5,
                        )
                        try:
                            docx_to_pdf_via_libreoffice(source, output, progress)
                        except Exception as fallback_error:
                            raise RuntimeError(
                                f"Word 原生转换失败：{native_error}\n"
                                f"LibreOffice 回退也失败：{fallback_error}"
                            ) from fallback_error
                elif (
                    not self.is_macos
                    and _engine_result_selectable(word_com_available())
                ):
                    self._convert_with_windows_fallback(
                        "Word",
                        docx_to_pdf_via_word,
                        docx_to_pdf_via_libreoffice,
                        source,
                        output,
                        progress,
                    )
                elif _engine_result_selectable(libreoffice_available()):
                    docx_to_pdf_via_libreoffice(source, output, progress)
                else:
                    raise RuntimeError("Word → PDF 需要 Microsoft Word 或 LibreOffice")

            return convert_docx

        if spec.key == "powerpoint_to_pdf":
            def convert_presentation(source, output, progress):
                if self.is_macos and self._method_availability(POWERPOINT_NATIVE):
                    try:
                        self._convert_with_macos_backend(
                            POWERPOINT_NATIVE,
                            self._mac_powerpoint_backend,
                            source,
                            output,
                            progress,
                        )
                    except MacOSOfficeError as native_error:
                        if not _engine_result_selectable(libreoffice_available()):
                            raise
                        output_path = Path(output)
                        if output_path.exists():
                            output_path.unlink()
                        progress(
                            "PowerPoint 原生转换失败，改用 LibreOffice 重试（详见自动化权限状态）...",
                            5,
                        )
                        try:
                            presentation_to_pdf_via_libreoffice(source, output, progress)
                        except Exception as fallback_error:
                            raise RuntimeError(
                                f"PowerPoint 原生转换失败：{native_error}\n"
                                f"LibreOffice 回退也失败：{fallback_error}"
                            ) from fallback_error
                elif (
                    not self.is_macos
                    and _engine_result_selectable(powerpoint_com_available())
                ):
                    self._convert_with_windows_fallback(
                        "PowerPoint",
                        presentation_to_pdf_via_powerpoint,
                        presentation_to_pdf_via_libreoffice,
                        source,
                        output,
                        progress,
                    )
                elif _engine_result_selectable(libreoffice_available()):
                    presentation_to_pdf_via_libreoffice(source, output, progress)
                else:
                    raise RuntimeError(
                        "PowerPoint → PDF 需要 Microsoft PowerPoint 或 LibreOffice"
                    )

            return convert_presentation

        raise RuntimeError(f"未知转换方向: {spec.key}")

    def _progress_update(self, current, total, message, pct):
        self.root.after(0, self._set_progress, current, total, message, pct)

    def _set_progress(self, current, total, message, pct):
        self.progress["value"] = pct
        self.progress_text_var.set(f"[{current}/{total}] {message} · 总进度 {pct}%")
        self.log(f"[{current}/{total}] {message}")

    def _status_update(self, result: BatchResult):
        self.root.after(0, self._apply_status_update, result)

    def _apply_status_update(self, result: BatchResult):
        item_id = next(
            (
                current_id for current_id, path in self._tree_paths.items()
                if self._path_key(path) == self._path_key(result.source)
            ),
            None,
        )
        if item_id:
            values = list(self.file_tree.item(item_id, "values"))
            values[2] = self.STATUS_TEXT[result.status]
            values[3] = str(result.output) if result.output and result.status != "cancelled" else ""
            self.file_tree.item(item_id, values=values, tags=(result.status,))
            self.file_tree.see(item_id)
            item_number = self.file_tree.index(item_id) + 1
        else:
            item_number = 1

        total = max(1, len(self.input_paths))
        prefix = f"[{item_number}/{total}]"

        if result.status == "running":
            self.log(f"{prefix} 开始: {result.source.name}")
        elif result.status == "success":
            self.log(f"{prefix} 成功: {result.output}")
        elif result.status == "failed":
            self.log(f"{prefix} 失败: {result.source.name} - {result.error}")
        elif result.status == "cancelled":
            self.log(f"{prefix} 已取消: {result.source.name}")

    def _stop_cancel_wait_timer(self):
        after_id = getattr(self, "_cancel_wait_after_id", None)
        if after_id is not None:
            after_cancel = getattr(self.root, "after_cancel", None)
            if callable(after_cancel):
                try:
                    after_cancel(after_id)
                except (RuntimeError, tk.TclError):
                    pass
        self._cancel_wait_after_id = None
        self._cancel_wait_started_at = None

    def _cancel_wait_message(self, elapsed_seconds):
        if getattr(self, "_close_after_batch", False):
            action = "当前文件结束后安全退出"
        else:
            action = "当前文件处理结束后取消"
        return f"正在等待{action}...（已等待 {elapsed_seconds} 秒）"

    def _cancel_wait_tick(self):
        self._cancel_wait_after_id = None
        if not self._is_converting or not self._cancel_requested():
            self._stop_cancel_wait_timer()
            return
        started_at = getattr(self, "_cancel_wait_started_at", None)
        if started_at is None:
            return
        elapsed_seconds = max(0, int(time.monotonic() - started_at))
        self.progress_text_var.set(self._cancel_wait_message(elapsed_seconds))
        self._cancel_wait_after_id = self.root.after(
            1000,
            self._cancel_wait_tick,
        )

    def _start_cancel_wait_timer(self):
        self._stop_cancel_wait_timer()
        self._cancel_wait_started_at = time.monotonic()
        self.progress_text_var.set(self._cancel_wait_message(0))
        self._cancel_wait_after_id = self.root.after(
            1000,
            self._cancel_wait_tick,
        )

    def cancel_conversion(self):
        if not self._is_converting or self._cancel_event.is_set():
            return
        self._cancel_event.set()
        self.cancel_btn.config(state="disabled")
        self._start_cancel_wait_timer()
        self.log("已请求取消；当前文件结束后将停止后续任务")

    def _on_batch_complete(self, results, warnings=None):
        self._stop_cancel_wait_timer()
        self._set_busy(False)
        successful = [result for result in results if result.status == "success"]
        failed = [result for result in results if result.status == "failed"]
        cancelled = [result for result in results if result.status == "cancelled"]
        self.output_paths = [result.output for result in successful if result.output]

        summary = f"成功 {len(successful)}，失败 {len(failed)}，取消 {len(cancelled)}"
        self.queue_summary_var.set(f"批次完成 · {summary}")
        self.progress_text_var.set(summary)
        if not cancelled:
            self.progress["value"] = 100
        self.log(f"批次结束: {summary}")
        if self._close_after_batch:
            self.root.destroy()
            return

        details = ""
        if failed:
            shown = failed[:3]
            details = "\n\n失败详情:\n" + "\n".join(
                f"- {result.source.name}: {result.error}" for result in shown
            )
            if len(failed) > len(shown):
                details += f"\n- 另有 {len(failed) - len(shown)} 个失败，请查看日志"

        warnings = tuple(warnings or ())
        if warnings:
            shown_warnings = warnings[:5]
            details += "\n\n转换质量提醒:\n" + "\n".join(
                f"- {warning}" for warning in shown_warnings
            )
            if len(warnings) > len(shown_warnings):
                details += (
                    f"\n- 另有 {len(warnings) - len(shown_warnings)} 个提醒，"
                    "请查看日志"
                )
            for warning in warnings:
                self.log(f"转换质量提醒: {warning}")

        message = summary + details
        if warnings:
            messagebox.showwarning("转换完成，请仔细核对", message)
        elif len(successful) == len(results):
            messagebox.showinfo("转换完成", message)
        elif not successful and failed and not cancelled:
            messagebox.showerror("转换失败", message)
        else:
            messagebox.showwarning("批次已结束", message)

    def _on_batch_error(self, error):
        self._stop_cancel_wait_timer()
        self._set_busy(False)
        self.progress_text_var.set("批次异常终止")
        self.log(f"批次异常终止: {error}")
        if self._close_after_batch:
            self.root.destroy()
            return
        messagebox.showerror("转换错误", error)

    def _on_window_close(self):
        if self._is_converting:
            confirmed = messagebox.askyesno(
                "安全退出",
                "当前文件仍在转换。是否在当前文件完成或失败后退出？\n\n"
                "程序不会强制终止 Office、LibreOffice 或 AppleScript。",
            )
            if not confirmed:
                return
            self._close_after_batch = True
            self.cancel_conversion()
            self.log("窗口将在当前文件结束后关闭")
            return
        self.root.destroy()

    def _set_busy(self, busy):
        self._is_converting = busy
        self._update_action_states()

    def _update_action_states(self):
        editable_state = "disabled" if self._is_converting or self._is_installing else "normal"
        self.select_pdf_btn.config(state=editable_state)
        self.select_docx_btn.config(state=editable_state)
        self.select_ppt_btn.config(state=editable_state)
        self.clear_btn.config(
            state="normal" if self.input_paths and editable_state == "normal" else "disabled"
        )
        self.remove_btn.config(
            state="normal" if self.file_tree.selection() and editable_state == "normal" else "disabled"
        )
        self._refresh_target_zones()
        self.source_output_radio.config(state=editable_state)
        self.custom_output_radio.config(state=editable_state)
        browse_enabled = (
            not self._is_converting and not self._is_installing
            and self.output_mode_var.get() == "custom"
        )
        self.browse_output_btn.config(state="normal" if browse_enabled else "disabled")
        can_start = (
            self.input_paths and self._conversion_detection_ready()
            and not self._is_converting and not self._is_installing
        )
        self.convert_btn.config(state="normal" if can_start else "disabled")
        self.cancel_btn.config(state="normal" if self._is_converting else "disabled")

        spec = self._current_spec()
        for value, radio in self._method_widgets.items():
            available = self._method_availability(value)
            selectable = available is True or (
                available is False and value in (WORD_NATIVE, "libreoffice")
            )
            enabled = (
                not self._is_converting and not self._is_installing
                and (spec is None or spec.key == "pdf_to_word")
                and selectable
            )
            radio.config(state="normal" if enabled else "disabled")

    def open_output_folder(self):
        target = None
        if self.output_paths:
            target = self.output_paths[0].parent
        elif self.output_mode_var.get() == "custom" and self.output_dir_var.get():
            candidate = Path(self.output_dir_var.get())
            if candidate.is_dir():
                target = candidate
        elif self.input_paths:
            target = self.input_paths[0].parent

        if target and target.is_dir():
            try:
                self.platform_services.open_directory(target)
            except Exception as exc:
                logging.exception("Could not open output directory")
                messagebox.showerror("无法打开目录", str(exc))
        else:
            messagebox.showinfo("提示", "当前没有可打开的输出目录")


def main():
    root = None
    if TkinterDnD is not None:
        try:
            root = TkinterDnD.Tk()
            root._pdf_converter_dnd_available = True
        except Exception:
            logging.exception("TkDND root creation failed; using standard Tk")
    if root is None:
        root = tk.Tk()
        root._pdf_converter_dnd_available = False
    try:
        root.iconbitmap(default="")
    except Exception:
        pass
    ConverterApp(root)
    root.mainloop()


if __name__ == "__main__":
    main()
