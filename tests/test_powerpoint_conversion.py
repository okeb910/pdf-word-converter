import sys
import tempfile
import types
import unittest
from pathlib import Path
from unittest.mock import Mock, patch

import pymupdf
from pptx import Presentation

import pdf_word_converter as converter


class PdfToPowerPointTests(unittest.TestCase):
    def _create_mixed_page_pdf(self, path):
        document = pymupdf.open()
        portrait = document.new_page(width=595, height=842)
        portrait.draw_rect(pymupdf.Rect(40, 40, 555, 802), color=(0.8, 0.1, 0.1))
        portrait.insert_text((72, 100), "Portrait page", fontsize=24)
        landscape = document.new_page(width=842, height=595)
        landscape.draw_rect(pymupdf.Rect(40, 40, 802, 555), color=(0.1, 0.3, 0.8))
        landscape.insert_text((72, 100), "Landscape page", fontsize=24)
        document.save(path)
        document.close()

    def test_adaptive_dpi_caps_large_pages(self):
        self.assertEqual(converter.calculate_adaptive_dpi(595, 842), 200)
        self.assertEqual(converter.calculate_adaptive_dpi(7200, 3600), 32)
        with self.assertRaises(ValueError):
            converter.calculate_adaptive_dpi(0, 100)

    def test_mixed_page_pdf_creates_centered_uncropped_slides(self):
        with tempfile.TemporaryDirectory() as tmp:
            folder = Path(tmp)
            source = folder / "中文 混合页面.pdf"
            output = folder / "中文 混合页面.pptx"
            self._create_mixed_page_pdf(source)
            progress = []

            converter.pdf_to_pptx_via_images(
                str(source), str(output), lambda message, pct: progress.append((message, pct))
            )

            presentation = Presentation(output)
            self.assertEqual(len(presentation.slides), 2)
            self.assertAlmostEqual(
                presentation.slide_width / presentation.slide_height,
                595 / 842,
                places=3,
            )
            for slide in presentation.slides:
                self.assertEqual(len(slide.shapes), 1)
                picture = slide.shapes[0]
                self.assertGreater(picture.width, 0)
                self.assertGreater(picture.height, 0)
                self.assertGreaterEqual(picture.left, 0)
                self.assertGreaterEqual(picture.top, 0)
                self.assertLessEqual(picture.left + picture.width, presentation.slide_width)
                self.assertLessEqual(picture.top + picture.height, presentation.slide_height)
            self.assertGreater(presentation.slides[1].shapes[0].top, 0)
            self.assertEqual(progress[-1][1], 100)


class PowerPointExportTests(unittest.TestCase):
    def test_powerpoint_com_uses_read_only_hidden_open_and_cleans_up(self):
        comtypes_module = types.ModuleType("comtypes")
        client_module = types.ModuleType("comtypes.client")
        comtypes_module.CoInitialize = Mock()
        comtypes_module.CoUninitialize = Mock()
        comtypes_module.client = client_module

        presentation = Mock()
        powerpoint = Mock()
        powerpoint.Presentations.Open.return_value = presentation
        client_module.CreateObject = Mock(return_value=powerpoint)

        with tempfile.TemporaryDirectory() as tmp:
            source = Path(tmp) / "deck.pptx"
            output = Path(tmp) / "deck.pdf"
            source.touch()
            with patch.dict(
                sys.modules,
                {"comtypes": comtypes_module, "comtypes.client": client_module},
            ):
                converter.presentation_to_pdf_via_powerpoint(
                    str(source), str(output), lambda *_args: None
                )

        client_module.CreateObject.assert_called_once_with(
            "PowerPoint.Application", dynamic=True
        )
        powerpoint.Presentations.Open.assert_called_once_with(
            str(source.absolute()), True, False, False
        )
        presentation.SaveAs.assert_called_once_with(str(output.absolute()), 32)
        presentation.Close.assert_called_once_with()
        powerpoint.Quit.assert_called_once_with()
        comtypes_module.CoUninitialize.assert_called_once_with()

    def test_libreoffice_moves_isolated_powerpoint_output(self):
        with tempfile.TemporaryDirectory() as tmp:
            folder = Path(tmp)
            source = folder / "slides.ppt"
            output = folder / "renamed.pdf"
            source.touch()

            def fake_run(command, **_kwargs):
                out_dir = Path(command[command.index("--outdir") + 1])
                (out_dir / "slides.pdf").write_bytes(b"pdf")
                return types.SimpleNamespace(returncode=0, stdout="", stderr="")

            with patch.object(converter, "_get_lo_path", return_value="soffice"), patch.object(
                converter.subprocess, "run", side_effect=fake_run
            ):
                converter.presentation_to_pdf_via_libreoffice(
                    str(source), str(output), lambda *_args: None
                )

            self.assertEqual(output.read_bytes(), b"pdf")


if __name__ == "__main__":
    unittest.main()
